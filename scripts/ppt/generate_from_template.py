#!/usr/bin/env python3
"""
PPT Template Fill — 기존 PPTX 레이아웃 유지 + 텍스트 교체

기존 PPTX 파일의 레이아웃/디자인을 그대로 유지하면서
텍스트 내용만 교체합니다. python-pptx XML 노드 교체 방식.

Usage:
    python generate_from_template.py template.pptx replacements.json
    python generate_from_template.py template.pptx replacements.json --output out.pptx
    python generate_from_template.py template.pptx replacements.json --dry-run
    python generate_from_template.py template.pptx --inspect   # 슬라이드/쉐이프 목록 출력

replacements.json 형식 (두 가지 지원):

형식 A — 검색·교체 (위치 무관):
[
    {"search": "OLD TEXT", "replace": "NEW TEXT"},
    {"search": "제목 텍스트", "replace": "새 제목", "slide": 0}
]

형식 B — 쉐이프 직접 지정:
{
    "slides": {
        "0": {"제목 1": "새 제목", "내용 플레이스홀더 2": "새 내용"},
        "1": {"Shape Name": "New text"}
    }
}

제약:
- 텍스트 영역만 교체 가능 (이미지/SmartArt 불가)
- 폰트/색상/크기는 첫 번째 Run에서 상속
- 줄바꿈: \\n → 단락 분리

Requirements:
    pip install python-pptx
"""

import argparse
import copy
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn


# ─── Helpers ──────────────────────────────────────────────────────

def get_shape_text(shape):
    """쉐이프의 전체 텍스트 반환."""
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def set_text_preserving_format(text_frame, new_text: str):
    """
    텍스트프레임의 텍스트를 교체하되 포맷(폰트·색·크기)은 유지.

    전략:
    1. 첫 번째 단락의 첫 번째 Run에서 포맷 상속
    2. \\n으로 분리된 라인마다 새 단락 생성
    3. 기존 단락/Run 전부 제거 후 재구성
    """
    tf = text_frame

    # 참조용 포맷 추출 (첫 Run에서)
    ref_run_elem = None
    ref_para_elem = None
    for para in tf.paragraphs:
        ref_para_elem = para._p
        for run in para.runs:
            ref_run_elem = run._r
            break
        if ref_run_elem is not None:
            break

    # 기존 단락 제거 (XML 직접 조작)
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)

    lines = new_text.split('\n') if new_text else ['']
    for line in lines:
        # 새 단락 XML 생성
        if ref_para_elem is not None:
            new_para = copy.deepcopy(ref_para_elem)
            # 기존 Run 전부 제거
            for r in new_para.findall(qn('a:r')):
                new_para.remove(r)
            # 기존 줄바꿈(br) 제거
            for br in new_para.findall(qn('a:br')):
                new_para.remove(br)
        else:
            new_para = _make_empty_para()

        # Run 추가
        if ref_run_elem is not None:
            new_run = copy.deepcopy(ref_run_elem)
            # 텍스트 설정
            t_elem = new_run.find(qn('a:t'))
            if t_elem is None:
                from lxml import etree
                t_elem = etree.SubElement(new_run, qn('a:t'))
            t_elem.text = line
        else:
            new_run = _make_run(line)

        new_para.append(new_run)
        txBody.append(new_para)


def _make_empty_para():
    from lxml import etree
    return etree.fromstring('<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')


def _make_run(text: str):
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    r = etree.fromstring(f'<a:r xmlns:a="{ns}"><a:t/></a:r>')
    r.find(qn('a:t')).text = text
    return r


# ─── Inspect mode ─────────────────────────────────────────────────

def inspect_presentation(prs: Presentation):
    """슬라이드·쉐이프·텍스트 목록 출력."""
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n[Slide {slide_idx}]")
        for shape in slide.shapes:
            text = get_shape_text(shape)
            text_preview = repr(text[:60]) if text else "(no text)"
            print(f"  shape_id={shape.shape_id} name={repr(shape.name)} type={shape.shape_type}")
            if text:
                print(f"    text={text_preview}")


# ─── Replace: Format A (search/replace) ───────────────────────────

def apply_search_replace(prs: Presentation, replacements: list, dry_run: bool) -> int:
    changes = 0
    for rep in replacements:
        search = rep.get("search", "")
        replace = rep.get("replace", "")
        target_slide = rep.get("slide", None)  # None = all slides

        if not search:
            continue

        slides = prs.slides
        if target_slide is not None:
            slides = [prs.slides[int(target_slide)]]

        for slide in slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                current = get_shape_text(shape)
                if search in current:
                    new_text = current.replace(search, replace)
                    print(f"  {'[DRY]' if dry_run else '[OK]'} slide={prs.slides.index(slide)} "
                          f"shape={repr(shape.name)}: {repr(search[:30])} → {repr(replace[:30])}")
                    if not dry_run:
                        set_text_preserving_format(shape.text_frame, new_text)
                    changes += 1
    return changes


# ─── Replace: Format B (shape-direct) ─────────────────────────────

def apply_shape_direct(prs: Presentation, slides_map: dict, dry_run: bool) -> int:
    changes = 0
    for slide_idx_str, shape_map in slides_map.items():
        slide_idx = int(slide_idx_str)
        if slide_idx >= len(prs.slides):
            print(f"  [WARN] slide {slide_idx} out of range (total={len(prs.slides)})")
            continue
        slide = prs.slides[slide_idx]

        for shape_name, new_text in shape_map.items():
            matched = False
            for shape in slide.shapes:
                if shape.name == shape_name and shape.has_text_frame:
                    old_text = get_shape_text(shape)
                    print(f"  {'[DRY]' if dry_run else '[OK]'} slide={slide_idx} "
                          f"shape={repr(shape_name)}: {repr(old_text[:30])} → {repr(new_text[:30])}")
                    if not dry_run:
                        set_text_preserving_format(shape.text_frame, new_text)
                    changes += 1
                    matched = True
                    break
            if not matched:
                print(f"  [WARN] slide={slide_idx} shape={repr(shape_name)} not found")
    return changes


# ─── Main ─────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="PPT Template Fill — 텍스트 교체 엔진")
    parser.add_argument("template", help="입력 PPTX 경로")
    parser.add_argument("replacements", nargs="?", help="replacements.json 경로")
    parser.add_argument("--output", "-o", help="출력 PPTX 경로 (기본: template_filled.pptx)")
    parser.add_argument("--dry-run", action="store_true", help="변경 없이 미리보기만")
    parser.add_argument("--inspect", action="store_true", help="슬라이드·쉐이프 목록 출력 후 종료")
    args = parser.parse_args()

    template_path = Path(args.template)
    if not template_path.exists():
        print(f"ERROR: template not found: {template_path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(template_path))

    if args.inspect:
        inspect_presentation(prs)
        return

    if not args.replacements:
        parser.error("replacements.json 필요 (--inspect 없이 실행 시)")

    rep_path = Path(args.replacements)
    if not rep_path.exists():
        print(f"ERROR: replacements not found: {rep_path}", file=sys.stderr)
        sys.exit(1)

    with open(rep_path, encoding="utf-8") as f:
        data = json.load(f)

    output_path = Path(args.output) if args.output else template_path.with_stem(template_path.stem + "_filled")

    print(f"Template: {template_path} ({len(prs.slides)} slides)")
    print(f"Output:   {output_path}")
    if args.dry_run:
        print("Mode:     DRY RUN (no changes written)")
    print()

    total_changes = 0

    if isinstance(data, list):
        # Format A
        total_changes = apply_search_replace(prs, data, args.dry_run)
    elif isinstance(data, dict) and "slides" in data:
        # Format B
        total_changes = apply_shape_direct(prs, data["slides"], args.dry_run)
    elif isinstance(data, dict):
        # Format B shorthand (top-level keys = slide indices)
        total_changes = apply_shape_direct(prs, data, args.dry_run)
    else:
        print("ERROR: replacements.json 형식 불인식", file=sys.stderr)
        sys.exit(1)

    print(f"\n총 {total_changes}건 {'미리보기' if args.dry_run else '교체 완료'}")

    if not args.dry_run:
        prs.save(str(output_path))
        print(f"저장: {output_path}")
        # Run QA if available
        qa_script = Path(__file__).parent / "qa-pptx.py"
        if qa_script.exists():
            import subprocess
            result = subprocess.run(
                [sys.executable, str(qa_script), str(output_path)],
                capture_output=True, text=True
            )
            if result.returncode != 0:
                print("\n[QA 경고]")
                print(result.stdout[-500:] if result.stdout else result.stderr[-200:])


if __name__ == "__main__":
    main()
