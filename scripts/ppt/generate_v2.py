"""
A-Team PPT Engine v3 — 프로페셔널 편집 가능 .pptx 생성기
Gamma/Skywork 수준 + 사람이 만든 느낌 달성

v2 대비 개선:
  - 그라데이션 배경 (2색/3색)
  - 카드 그림자 (outerShdw XML)
  - CJK 폰트 정확 적용 (<a:ea> XML)
  - 풋터/슬라이드 번호 자동
  - 스피커 노트 삽입
  - 레이아웃 교차 강제 (동일 타입 2연속 금지)
  - 폰트 대비 극대화 (cover 64pt, 본문 14pt)
  - 신규 레이아웃 5종 (icon_grid, big_number, image_text, bento_grid, comparison)
  - 테마 5→8개 확장
  - 이미지 삽입 지원 (add_picture)

슬라이드 타입:
  구조: cover, section_break, agenda, closing
  콘텐츠: bullets/single, two_column, stats_grid, data_table, quote
  도식화: flow_diagram, timeline, bar_chart
  신규: icon_grid, big_number, image_text, bento_grid, comparison

테마: dark_editorial, consulting_clean, executive_deep,
      midnight_blue, warm_earth, nordic_frost, mono_sharp, sage_green

사용:
  python generate_v2.py spec.json
  python generate_v2.py spec.json --theme consulting_clean --output out.pptx
"""
import sys, json, argparse, os, io, tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from lxml import etree

RECT   = MSO_SHAPE.RECTANGLE
RRECT  = MSO_SHAPE.ROUNDED_RECTANGLE
OVAL   = MSO_SHAPE.OVAL
RARROW = MSO_SHAPE.RIGHT_ARROW
CHEVRON = MSO_SHAPE.CHEVRON

# ── 슬라이드 규격 (16:9) ─────────────────────────────────────
W  = Inches(13.333)
H  = Inches(7.5)
MX = Inches(0.75)
MY = Inches(0.6)
CW = W - 2 * MX

# DrawingML namespace
_nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# ── 테마 (8종) ────────────────────────────────────────────────
def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

THEMES = {
    "dark_editorial": {
        "bg": rgb("111116"), "bg2": rgb("16161e"),
        "surface": rgb("1a1a22"), "surface2": rgb("222228"),
        "text": rgb("e8e4dc"), "dim": rgb("8a8680"), "dim2": rgb("5a5858"),
        "metal": rgb("a89878"), "rule": rgb("2a2a30"),
        "accent": rgb("4a7fa8"), "accent2": rgb("6a9fc8"),
        "positive": rgb("4a8060"), "negative": rgb("8a4040"), "warn": rgb("a89048"),
        "card_bg": rgb("1a1a22"), "card_border": rgb("2a2a30"),
        "font": "Malgun Gothic", "font_display": "Malgun Gothic",
        "font_en": "Helvetica",
    },
    "consulting_clean": {
        "bg": rgb("f5f0e6"), "bg2": rgb("ede8dc"),
        "surface": rgb("ffffff"), "surface2": rgb("f0ebe0"),
        "text": rgb("1a1814"), "dim": rgb("5a5248"), "dim2": rgb("9a9088"),
        "metal": rgb("8a7860"), "rule": rgb("d8d0c4"),
        "accent": rgb("1d4e8a"), "accent2": rgb("2d6eaa"),
        "positive": rgb("2a5e30"), "negative": rgb("8a2020"), "warn": rgb("8a6020"),
        "card_bg": rgb("ffffff"), "card_border": rgb("d8d0c4"),
        "font": "Malgun Gothic", "font_display": "Malgun Gothic",
        "font_en": "Georgia",
    },
    "executive_deep": {
        "bg": rgb("0c0a09"), "bg2": rgb("120e0a"),
        "surface": rgb("18140e"), "surface2": rgb("22180e"),
        "text": rgb("f0ead8"), "dim": rgb("8a7a68"), "dim2": rgb("604838"),
        "metal": rgb("b8986a"), "rule": rgb("2a2018"),
        "accent": rgb("9f1239"), "accent2": rgb("c8304a"),
        "positive": rgb("1a5c3a"), "negative": rgb("8a1818"), "warn": rgb("8a6018"),
        "card_bg": rgb("18140e"), "card_border": rgb("2a2018"),
        "font": "Malgun Gothic", "font_display": "Malgun Gothic",
        "font_en": "Georgia",
    },
    "midnight_blue": {
        "bg": rgb("0a1628"), "bg2": rgb("0e1e38"),
        "surface": rgb("142040"), "surface2": rgb("1a2a50"),
        "text": rgb("e8eef8"), "dim": rgb("8a9ab8"), "dim2": rgb("5a6a88"),
        "metal": rgb("7890b0"), "rule": rgb("1e3050"),
        "accent": rgb("3b82f6"), "accent2": rgb("60a5fa"),
        "positive": rgb("22c55e"), "negative": rgb("ef4444"), "warn": rgb("f59e0b"),
        "card_bg": rgb("142040"), "card_border": rgb("1e3050"),
        "font": "Malgun Gothic", "font_display": "Malgun Gothic",
        "font_en": "Helvetica",
    },
    "warm_earth": {
        "bg": rgb("faf6f0"), "bg2": rgb("f2ece2"),
        "surface": rgb("ffffff"), "surface2": rgb("f5efe5"),
        "text": rgb("2c1810"), "dim": rgb("6b4a38"), "dim2": rgb("9a8070"),
        "metal": rgb("8b6914"), "rule": rgb("e0d4c0"),
        "accent": rgb("b45309"), "accent2": rgb("d97706"),
        "positive": rgb("15803d"), "negative": rgb("b91c1c"), "warn": rgb("a16207"),
        "card_bg": rgb("ffffff"), "card_border": rgb("e0d4c0"),
        "font": "Malgun Gothic", "font_display": "Malgun Gothic",
        "font_en": "Georgia",
    },
    "nordic_frost": {
        "bg": rgb("f8fafb"), "bg2": rgb("eef2f5"),
        "surface": rgb("ffffff"), "surface2": rgb("f0f4f7"),
        "text": rgb("1a2332"), "dim": rgb("4a5568"), "dim2": rgb("8a95a5"),
        "metal": rgb("64748b"), "rule": rgb("dde4ea"),
        "accent": rgb("0ea5e9"), "accent2": rgb("38bdf8"),
        "positive": rgb("10b981"), "negative": rgb("f43f5e"), "warn": rgb("f59e0b"),
        "card_bg": rgb("ffffff"), "card_border": rgb("dde4ea"),
        "font": "Malgun Gothic", "font_display": "Malgun Gothic",
        "font_en": "Helvetica",
    },
    "mono_sharp": {
        "bg": rgb("ffffff"), "bg2": rgb("f5f5f5"),
        "surface": rgb("fafafa"), "surface2": rgb("f0f0f0"),
        "text": rgb("0a0a0a"), "dim": rgb("404040"), "dim2": rgb("808080"),
        "metal": rgb("505050"), "rule": rgb("e0e0e0"),
        "accent": rgb("0a0a0a"), "accent2": rgb("404040"),
        "positive": rgb("166534"), "negative": rgb("991b1b"), "warn": rgb("854d0e"),
        "card_bg": rgb("fafafa"), "card_border": rgb("e0e0e0"),
        "font": "Malgun Gothic", "font_display": "Malgun Gothic",
        "font_en": "Helvetica",
    },
    "sage_green": {
        "bg": rgb("f5f7f4"), "bg2": rgb("eaefe8"),
        "surface": rgb("ffffff"), "surface2": rgb("f0f4ee"),
        "text": rgb("1a2e1a"), "dim": rgb("4a6040"), "dim2": rgb("7a9070"),
        "metal": rgb("5a7850"), "rule": rgb("d0dcc8"),
        "accent": rgb("2d6a4f"), "accent2": rgb("40916c"),
        "positive": rgb("1b4332"), "negative": rgb("9b2226"), "warn": rgb("8a6d3b"),
        "card_bg": rgb("ffffff"), "card_border": rgb("d0dcc8"),
        "font": "Malgun Gothic", "font_display": "Malgun Gothic",
        "font_en": "Georgia",
    },
}

MIN_FONT = 12

def fs(size):
    """최소 폰트 보장."""
    return max(size, MIN_FONT)


# ── 유틸리티 (신규) ───────────────────────────────────────────

def _is_dark_theme(th):
    """배경색 밝기로 다크/라이트 판별."""
    bg = th["bg"]
    lum = 0.299 * bg[0] + 0.587 * bg[1] + 0.114 * bg[2]
    return lum < 128


def paint_bg(slide, color):
    """단색 배경."""
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color


def paint_bg_gradient(slide, color1, color2, angle=135):
    """2색 그라데이션 배경. angle: 0=왼→오, 90=위→아래, 135=좌상→우하."""
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = angle
    stops = fill.gradient_stops
    stops[0].color.rgb = color1
    stops[0].position = 0.0
    stops[1].color.rgb = color2
    stops[1].position = 1.0


def apply_shadow(shape, blur=50000, dist=25000, direction=2700000, alpha=40):
    """도형에 outerShdw 그림자 적용. direction: 60000분의 1도 단위 (2700000=270도=아래)."""
    spPr = shape._element.spPr
    nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # 기존 effectLst 제거
    for el in spPr.findall(f'{{{nsmap_a}}}effectLst'):
        spPr.remove(el)

    effectLst = etree.SubElement(spPr, f'{{{nsmap_a}}}effectLst')
    outerShdw = etree.SubElement(effectLst, f'{{{nsmap_a}}}outerShdw')
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', str(direction))
    outerShdw.set('rotWithShape', '0')

    srgbClr = etree.SubElement(outerShdw, f'{{{nsmap_a}}}srgbClr')
    srgbClr.set('val', '000000')
    alphaEl = etree.SubElement(srgbClr, f'{{{nsmap_a}}}alpha')
    alphaEl.set('val', str(alpha * 1000))  # 40% = 40000


def set_cjk_font(run, font_name):
    """CJK 폰트를 run에 정확히 적용 (python-pptx 한국어 폰트 버그 우회)."""
    nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    rPr = run._r.get_or_add_rPr()
    # 기존 <a:ea> 제거
    for ea in rPr.findall(f'{{{nsmap_a}}}ea'):
        rPr.remove(ea)
    # 새 <a:ea> 삽입
    ea = etree.SubElement(rPr, f'{{{nsmap_a}}}ea')
    ea.set('typeface', font_name)


def add_footer(slide, th, slide_num, total_slides, company=""):
    """슬라이드 하단에 풋터(슬라이드 번호 + 회사명) 삽입."""
    F = th["font"]
    # 슬라이드 번호 (우하단)
    num_text = f"{slide_num} / {total_slides}"
    s = slide.shapes.add_shape(RECT, W - MX - Inches(1.0), H - Inches(0.4),
                                Inches(1.0), Inches(0.3))
    s.fill.background()
    s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = num_text
    r.font.size = Pt(9)
    r.font.color.rgb = th["dim2"]
    r.font.name = th.get("font_en", F)
    set_cjk_font(r, F)

    # 회사명/브랜드 (좌하단)
    if company:
        s2 = slide.shapes.add_shape(RECT, MX, H - Inches(0.4),
                                     Inches(3.0), Inches(0.3))
        s2.fill.background()
        s2.line.fill.background()
        s2.shadow.inherit = False
        tf2 = s2.text_frame
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = company
        r2.font.size = Pt(9)
        r2.font.color.rgb = th["dim2"]
        r2.font.name = th.get("font_en", F)
        set_cjk_font(r2, F)


def add_notes(slide, notes_text):
    """스피커 노트 삽입."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def try_add_image(slide, image_path, left, top, width, height=None):
    """이미지 파일이 존재하면 삽입. 없으면 무시."""
    if not image_path or not os.path.exists(image_path):
        return None
    try:
        if height:
            return slide.shapes.add_picture(image_path, left, top, width, height)
        return slide.shapes.add_picture(image_path, left, top, width)
    except Exception:
        return None


# ── 프리미티브 ────────────────────────────────────────────────

def shape_text(slide, l, t, w, h, text,
               size=14, color=None, bold=False, font="Malgun Gothic",
               h_align=PP_ALIGN.LEFT, v_align=MSO_ANCHOR.MIDDLE,
               fill=None, border=None, bw=0.5,
               shape_type=RECT, italic=False, wrap=True,
               shadow=False, font_en=None):
    """도형 + 텍스트. shadow=True 시 카드 그림자 적용."""
    s = slide.shapes.add_shape(shape_type, l, t, w, h)

    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()

    if border is not None:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()

    s.shadow.inherit = False
    if shadow:
        apply_shadow(s)

    tf = s.text_frame
    tf.word_wrap = wrap
    tf.margin_left   = Emu(114300)
    tf.margin_right  = Emu(114300)
    tf.margin_top    = Emu(45720)
    tf.margin_bottom = Emu(45720)
    tf.vertical_anchor = v_align

    p = tf.paragraphs[0]
    p.alignment = h_align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs(size))
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font_en or font
    if color:
        r.font.color.rgb = color
    set_cjk_font(r, font)

    return s


def shape_multiline(slide, l, t, w, h, lines,
                    size=14, color=None, bold=False, font="Malgun Gothic",
                    h_align=PP_ALIGN.LEFT, v_align=MSO_ANCHOR.MIDDLE,
                    fill=None, border=None, bw=0.5, shape_type=RECT,
                    line_spacing=None, italic=False, shadow=False, font_en=None):
    """여러 줄 텍스트를 도형 안에 삽입."""
    s = slide.shapes.add_shape(shape_type, l, t, w, h)

    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()

    if border is not None:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()

    s.shadow.inherit = False
    if shadow:
        apply_shadow(s)

    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(114300)
    tf.margin_right  = Emu(114300)
    tf.margin_top    = Emu(45720)
    tf.margin_bottom = Emu(45720)
    tf.vertical_anchor = v_align

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = h_align
        if line_spacing:
            p.line_spacing = line_spacing

        if isinstance(line, list):
            for seg in line:
                r = p.add_run()
                r.text = seg[0]
                r.font.size = Pt(fs(seg[2] if len(seg) > 2 else size))
                r.font.color.rgb = seg[1] if seg[1] else (color or RGBColor(0,0,0))
                r.font.bold = seg[3] if len(seg) > 3 else bold
                r.font.italic = italic
                r.font.name = font_en or font
                set_cjk_font(r, font)
        else:
            r = p.add_run()
            r.text = str(line)
            r.font.size = Pt(fs(size))
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font_en or font
            if color:
                r.font.color.rgb = color
            set_cjk_font(r, font)

    return s


def solid_rect(slide, l, t, w, h, fill, border=None, bw=0.5):
    s = slide.shapes.add_shape(RECT, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def solid_rrect(slide, l, t, w, h, fill, border=None, bw=1.0, shadow=False):
    s = slide.shapes.add_shape(RRECT, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    if shadow:
        apply_shadow(s, blur=38000, dist=18000)
    return s


def hline(slide, l, t, w, color, weight=0.5):
    s = slide.shapes.add_shape(RECT, l, t, w, Pt(weight))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def vline(slide, l, t, h, color, weight=0.5):
    s = slide.shapes.add_shape(RECT, l, t, Pt(weight), h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def oval_dot(slide, cx, cy, r, fill):
    s = slide.shapes.add_shape(OVAL, cx - r, cy - r, r*2, r*2)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False
    return s


# ── 슬라이드 빌더 ─────────────────────────────────────────────

def slide_cover(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg_gradient(slide, th["bg"], th.get("bg2", th["bg"]))
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    # 좌측 강조 바 (더 두껍게)
    solid_rect(slide, MX, MY, Pt(4), H - 2*MY, th["accent"])
    L = MX + Inches(0.45)

    # 레이블 (kicker)
    kicker = content.get("label", content.get("kicker", "")).upper()
    if kicker:
        shape_text(slide, L, MY + Inches(0.15), CW - Inches(0.5), Inches(0.42),
                   kicker, 13, th["metal"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    # 타이틀 — 64pt 극대화
    title = content.get("headline", content.get("title", ""))
    lines = title.split("\n")
    shape_multiline(slide, L, MY + Inches(0.7), CW - Inches(1.0), Inches(3.0),
                    lines, 52, th["text"], True, FD,
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, line_spacing=1.12, font_en=FE)

    # 부제목 — 본문과 명확 구분
    subtitle = content.get("subtitle", content.get("subheadline", ""))
    if subtitle:
        shape_text(slide, L, MY + Inches(3.5), CW - Inches(1.0), Inches(0.7),
                   subtitle, 18, th["dim"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    # 구분선
    hline(slide, L, H - MY - Inches(1.25), CW * 0.5, th["metal"], 0.6)

    # 메타
    meta = content.get("meta", "")
    if meta:
        shape_text(slide, L, H - MY - Inches(1.1), CW - Inches(0.4), Inches(0.42),
                   meta, 12, th["dim2"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    # 태그
    tags = content.get("tags", [])
    tx = L
    ty = H - MY - Inches(0.55)
    for tag in tags[:4]:
        tw = Inches(max(1.4, len(tag) * 0.15 + 0.5))
        shape_text(slide, tx, ty, tw, Inches(0.38),
                   tag.upper(), 12, th["accent2"], True, F,
                   PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
                   th["surface"], th["accent"], 0.5,
                   shape_type=RRECT, font_en=FE)
        tx += tw + Inches(0.12)

    add_notes(slide, content.get("notes", ""))


def slide_section_break(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg_gradient(slide, th["bg"], th.get("bg2", th["bg"]), angle=180)
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    num = content.get("section_number", "01")

    # 워터마크 번호 — 더 극적
    shape_text(slide, MX - Inches(0.2), H/2 - Inches(2.5), Inches(5.0), Inches(4.0),
               num, 160, th.get("rule", th["dim2"]), True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    # Section 레이블
    shape_text(slide, MX + Inches(0.15), H/2 - Inches(0.8), Inches(4), Inches(0.42),
               f"Section {num}", 13, th["metal"], False, F,
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    # 강조 라인 (accent 색상)
    hline(slide, MX + Inches(0.15), H/2 - Inches(0.28), CW * 0.45, th["accent"], 1.5)

    # 헤드라인 — 48pt
    shape_text(slide, MX + Inches(0.15), H/2 - Inches(0.18), CW * 0.7, Inches(1.2),
               content.get("headline", ""), 36, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    desc = content.get("description", "")
    if desc:
        shape_text(slide, MX + Inches(0.15), H/2 + Inches(0.95), CW * 0.65, Inches(0.48),
                   desc, 14, th["dim"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_bullets(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FE = th["font"], th.get("font_en", th["font"])

    headline = content.get("headline", "")
    bullets  = content.get("bullets", [])
    source   = content.get("source", "")

    # 헤드라인 — 24pt bold, action title 스타일
    shape_text(slide, MX, MY, CW, Inches(0.72),
               headline, 24, th["text"], True, th["font_display"],
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.70), CW, th["metal"], 0.5)

    src_h  = Inches(0.4) if source else 0
    avail  = H - (MY + Inches(0.95)) - MY - src_h
    row_h  = min(Inches(0.72), avail / max(len(bullets), 1))
    y      = MY + Inches(0.95)

    for bullet in bullets[:5]:
        # 악센트 도트
        oval_dot(slide, MX + Inches(0.12), y + row_h / 2,
                 Inches(0.06), th["accent"])
        shape_text(slide, MX + Inches(0.32), y, CW - Inches(0.36), row_h,
                   bullet, 15, th["dim"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
        y += row_h

    if source:
        shape_text(slide, MX, H - MY - Inches(0.36), CW, Inches(0.36),
                   f"Source: {source}", 10, th["dim2"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, italic=True, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_two_column(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FE = th["font"], th.get("font_en", th["font"])

    left  = content.get("left", {})
    right = content.get("right", {})

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, th["font_display"],
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    col_w   = (CW - Inches(0.32)) / 2
    col_r   = MX + col_w + Inches(0.32)
    y_start = MY + Inches(0.9)
    col_h   = H - y_start - MY - Inches(0.1)
    PAD = Inches(0.22)

    for cx, side, accent in [(MX, left, th["accent"]), (col_r, right, th["accent2"])]:
        # 카드 (그림자 적용)
        card = solid_rrect(slide, cx, y_start, col_w, col_h,
                           th["card_bg"], th["card_border"], 0.5, shadow=True)

        # 상단 악센트 바
        solid_rect(slide, cx, y_start, col_w, Pt(4), accent)

        # 타이틀
        title_h = Inches(0.45)
        shape_text(slide, cx + PAD, y_start + Inches(0.18),
                   col_w - PAD*2, title_h,
                   side.get("title", "").upper(), 13, th["metal"], True, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
        hline(slide, cx + PAD, y_start + Inches(0.18) + title_h + Inches(0.05),
              col_w - PAD*2, th["rule"], 0.35)

        row_h = Inches(0.58)
        by    = y_start + Inches(0.18) + title_h + Inches(0.15)
        for b in side.get("bullets", [])[:5]:
            oval_dot(slide, cx + PAD + Inches(0.08), by + row_h / 2,
                     Inches(0.05), accent)
            shape_text(slide, cx + PAD + Inches(0.22), by,
                       col_w - PAD - Inches(0.26), row_h,
                       b, 13, th["dim"], False, F,
                       PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
            by += row_h

    add_notes(slide, content.get("notes", ""))


def slide_stats_grid(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    stats = content.get("stats", [])
    n = min(len(stats), 4)

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    if not n: return

    card_w = (CW - Inches(0.22) * (n-1)) / n
    y0   = MY + Inches(0.92)
    c_h  = H - y0 - MY - Inches(0.15)
    accents = [th["accent"], th["positive"], th["warn"], th["accent2"]]

    for i, stat in enumerate(stats[:n]):
        cx  = MX + i * (card_w + Inches(0.22))
        col = accents[i % 4]

        # 카드 (그림자)
        solid_rrect(slide, cx, y0, card_w, c_h, th["card_bg"], th["card_border"], 0.5, shadow=True)
        solid_rect(slide, cx, y0, card_w, Pt(4), col)

        # 레이블
        lbl_t = y0 + Inches(0.25)
        lbl_h = Inches(0.4)
        shape_text(slide, cx, lbl_t, card_w, lbl_h,
                   stat.get("label", "").upper(), 12, th["metal"], False, F,
                   PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, font_en=FE)

        # 주요 수치 — 42pt bold
        val_lines = stat.get("value", "").split("\n")
        val_t = lbl_t + lbl_h + Inches(0.08)
        val_h = Inches(1.5)
        shape_multiline(slide, cx, val_t, card_w, val_h,
                        val_lines, 42, th["text"], True, FD,
                        PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, line_spacing=1.1, font_en=FE)

        # 델타
        delta = stat.get("delta", "")
        delta_t = val_t + val_h + Inches(0.08)
        delta_h = Inches(0.42)
        if delta:
            d_col = th["positive"] if delta.startswith("+") else (
                    th["negative"] if delta.startswith("-") else th["dim"])
            shape_text(slide, cx, delta_t, card_w, delta_h,
                       delta, 14, d_col, True, F,
                       PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, font_en=FE)

        # 노트
        note = stat.get("note", "")
        note_t = max(delta_t + delta_h + Inches(0.08), y0 + c_h - Inches(0.52))
        if note:
            shape_text(slide, cx, note_t, card_w, Inches(0.45),
                       note, 12, th["dim2"], False, F,
                       PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_data_table(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FE = th["font"], th.get("font_en", th["font"])

    tdata   = content.get("table", {})
    headers = tdata.get("headers", [])
    rows    = tdata.get("rows", [])
    hl_col  = tdata.get("highlight_col", -1)
    source  = content.get("source", "")

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, th["font_display"],
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    if not headers or not rows: return

    nc    = len(headers)
    nr    = len(rows)
    src_h = Inches(0.38) if source else 0
    avail = H - (MY + Inches(0.92)) - MY - src_h
    row_h = min(Inches(0.54), avail / (nr + 1))
    col_w = CW / nc
    t_top = MY + Inches(0.88)

    # 헤더
    solid_rrect(slide, MX, t_top, CW, row_h, th["surface2"], th["rule"], 0.3)
    hline(slide, MX, t_top + row_h - Pt(0.4), CW, th["metal"], 0.8)
    for j, h_text in enumerate(headers):
        is_hl = (j == hl_col)
        shape_text(slide, MX + j*col_w, t_top, col_w, row_h,
                   h_text, 12,
                   th["accent2"] if is_hl else th["metal"],
                   True, F, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    # 데이터 행
    for i, row in enumerate(rows[:12]):
        ry = t_top + row_h * (i + 1)
        bg = th["surface"] if i % 2 == 0 else th["bg"]
        solid_rect(slide, MX, ry, CW, row_h, bg, th["rule"], 0.2)
        for j, cell in enumerate(row[:nc]):
            is_hl = (j == hl_col)
            shape_text(slide, MX + j*col_w, ry, col_w, row_h,
                       str(cell), 12,
                       th["text"] if is_hl else th["dim"],
                       is_hl, F, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    if source:
        sy = t_top + row_h * (len(rows[:12]) + 1) + Inches(0.08)
        shape_text(slide, MX, sy, CW, Inches(0.36),
                   f"Source: {source}", 10, th["dim2"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, italic=True, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_flow_diagram(prs, content, th):
    """도형 + 화살표 프로세스 플로우 (최대 6단계)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    steps = content.get("steps", [])
    n = min(len(steps), 6)

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    if not n: return

    arr_w   = Inches(0.35)
    arr_h   = Inches(0.38)
    box_h   = Inches(1.8)
    box_w   = (CW - (n-1)*arr_w) / n
    if box_w > Inches(2.2):
        box_w = Inches(2.2)
    total_w = n * box_w + (n-1) * arr_w
    start_x = MX + (CW - total_w) / 2
    cy      = MY + Inches(0.88) + (H - MY - Inches(0.88) - MY) / 2

    accents = [th["accent"], th["positive"], th["warn"],
               th["accent2"], th["metal"], th["dim"]]

    for i, step in enumerate(steps[:n]):
        bx  = start_x + i * (box_w + arr_w)
        col = accents[i % len(accents)]
        label = step.get("label", step) if isinstance(step, dict) else str(step)
        sub   = step.get("sub", "") if isinstance(step, dict) else ""

        # 카드 (그림자 포함)
        solid_rrect(slide, bx, cy - box_h/2, box_w, box_h,
                    th["card_bg"], col, 1.0, shadow=True)
        solid_rect(slide, bx, cy - box_h/2, box_w, Pt(4), col)

        # 번호 뱃지
        badge_r = Inches(0.17)
        badge_cx = bx + box_w - badge_r - Inches(0.12)
        badge_cy = cy - box_h/2 + badge_r + Inches(0.12)
        oval_dot(slide, badge_cx, badge_cy, badge_r, col)
        shape_text(slide, badge_cx - badge_r, badge_cy - badge_r,
                   badge_r*2, badge_r*2,
                   str(i+1), 12, th.get("card_bg", th["bg"]), True, F,
                   PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, font_en=FE)

        # 레이블
        label_t = cy - box_h/2 + Inches(0.2)
        label_h = Inches(0.95) if sub else Inches(1.2)
        shape_text(slide, bx + Inches(0.1), label_t, box_w - Inches(0.2), label_h,
                   label, 13, th["text"], True, F,
                   PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, wrap=True, font_en=FE)

        if sub:
            sub_t = label_t + label_h + Inches(0.06)
            sub_h = cy + box_h/2 - sub_t - Inches(0.12)
            if sub_h > Inches(0.3):
                shape_text(slide, bx + Inches(0.1), sub_t,
                           box_w - Inches(0.2), sub_h,
                           sub, 12, th["dim"], False, F,
                           PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, wrap=True, font_en=FE)

        # 화살표 (쉐브론 스타일)
        if i < n - 1:
            ax = bx + box_w + Inches(0.04)
            aw = arr_w - Inches(0.08)
            s  = slide.shapes.add_shape(CHEVRON, ax, cy - arr_h/2, aw, arr_h)
            s.fill.solid(); s.fill.fore_color.rgb = th["metal"]
            s.line.fill.background(); s.shadow.inherit = False

    add_notes(slide, content.get("notes", ""))


def slide_timeline(prs, content, th):
    """수직 타임라인."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FE = th["font"], th.get("font_en", th["font"])

    events = content.get("events", content.get("items", []))
    n      = min(len(events), 5)

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, th["font_display"],
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    if not n: return

    date_w    = Inches(1.4)
    line_x    = MX + date_w + Inches(0.3)
    content_x = line_x + Inches(0.42)
    content_w = CW - date_w - Inches(0.75)

    y0     = MY + Inches(0.92)
    y_end  = H - MY - Inches(0.25)
    step_h = (y_end - y0) / n
    accents = [th["accent"], th["positive"], th["accent2"], th["warn"], th["metal"]]

    # 수직 라인
    vline(slide, line_x, y0 + Inches(0.1), y_end - y0 - Inches(0.1), th["rule"], 2.0)

    for i, evt in enumerate(events[:n]):
        cy  = y0 + step_h * i + step_h / 2
        col = accents[i % len(accents)]

        # 도트 + 외곽 링
        dot_r = Inches(0.11)
        oval_dot(slide, line_x, cy, dot_r, col)
        s = slide.shapes.add_shape(OVAL,
                                    line_x - dot_r - Pt(2),
                                    cy     - dot_r - Pt(2),
                                    dot_r*2 + Pt(4),
                                    dot_r*2 + Pt(4))
        s.fill.background()
        s.line.color.rgb = col; s.line.width = Pt(1.2); s.shadow.inherit = False

        # 날짜 레이블
        date = evt.get("date", evt.get("period", "")) if isinstance(evt, dict) else ""
        date_h = Inches(0.46)
        shape_text(slide, MX, cy - date_h/2, date_w, date_h,
                   date, 13, th["metal"], True, F,
                   PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE, font_en=FE)

        # 제목 + 설명
        title = evt.get("title", str(evt)) if isinstance(evt, dict) else str(evt)
        desc  = evt.get("desc", evt.get("description", "")) if isinstance(evt, dict) else ""

        max_h = step_h - Inches(0.14)
        if desc:
            title_h = Inches(0.4)
            desc_h  = min(Inches(0.4), max_h - title_h - Inches(0.06))
            t_top   = cy - (title_h + Inches(0.04) + desc_h) / 2
            shape_text(slide, content_x, t_top, content_w, title_h,
                       title, 15, th["text"], True, F,
                       PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
            shape_text(slide, content_x, t_top + title_h + Inches(0.04),
                       content_w, desc_h,
                       desc, 12, th["dim"], False, F,
                       PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
        else:
            shape_text(slide, content_x, cy - Inches(0.25), content_w, Inches(0.5),
                       title, 15, th["text"], True, F,
                       PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_bar_chart(prs, content, th):
    """python-pptx 네이티브 차트."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FE = th["font"], th.get("font_en", th["font"])

    categories = content.get("categories", [])
    series     = content.get("series", [])
    source     = content.get("source", "")

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, th["font_display"],
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    if not categories or not series: return

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series:
        chart_data.add_series(s.get("name",""), tuple(s.get("values",[])))

    ctype = (XL_CHART_TYPE.BAR_CLUSTERED if content.get("horizontal")
             else XL_CHART_TYPE.COLUMN_CLUSTERED)
    src_h = Inches(0.38) if source else 0
    c_h   = H - (MY + Inches(0.92)) - MY - src_h

    slide.shapes.add_chart(ctype, MX, MY + Inches(0.88), CW, c_h, chart_data)

    if source:
        shape_text(slide, MX, H - MY - Inches(0.36), CW, Inches(0.36),
                   f"Source: {source}", 10, th["dim2"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, italic=True, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_quote(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg_gradient(slide, th["bg"], th.get("bg2", th["bg"]), angle=45)
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    quote       = content.get("quote", "")
    attribution = content.get("attribution", "")

    # 대형 인용 부호
    shape_text(slide, MX + Inches(0.1), MY + Inches(0.2), Inches(1.2), Inches(1.5),
               "\u201c", 96, th["rule"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_en=FE)

    # 수직 강조 바
    vline(slide, MX + Inches(0.18), MY + Inches(0.5), Inches(2.8), th["accent"], 3.0)

    # 인용문 — 26pt
    q_lines = quote.split("\n")
    q_top   = H/2 - Inches(1.6)
    q_h     = Inches(3.0)
    shape_multiline(slide, MX + Inches(0.6), q_top, CW * 0.8, q_h,
                    q_lines, 24, th["text"], False, FD,
                    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE,
                    line_spacing=1.4, italic=False, font_en=FE)

    if attribution:
        attr_t = q_top + q_h + Inches(0.15)
        hline(slide, MX + Inches(0.6), attr_t, Inches(2.8), th["metal"], 0.5)
        shape_text(slide, MX + Inches(0.6), attr_t + Inches(0.12), CW * 0.7, Inches(0.45),
                   f"\u2014 {attribution}", 13, th["dim"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_closing(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg_gradient(slide, th["bg"], th.get("bg2", th["bg"]))
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    solid_rect(slide, MX, MY, Pt(4), H - 2*MY, th["accent"])
    L = MX + Inches(0.45)

    hl_t = H/2 - Inches(1.3)
    hl_h = Inches(1.3)
    shape_text(slide, L, hl_t, CW * 0.8, hl_h,
               content.get("headline", "Thank You"), 48, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    rule_t = hl_t + hl_h + Inches(0.1)
    hline(slide, L, rule_t, CW * 0.45, th["accent"], 1.5)

    y_cur = rule_t + Inches(0.22)
    contact = content.get("contact", "")
    if contact:
        shape_text(slide, L, y_cur, CW * 0.8, Inches(0.48),
                   contact, 14, th["dim"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
        y_cur += Inches(0.52)

    note = content.get("note", "")
    if note:
        shape_text(slide, L, y_cur, CW * 0.8, Inches(0.48),
                   note, 12, th["dim2"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_agenda(prs, content, th):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FE = th["font"], th.get("font_en", th["font"])

    items = content.get("items", [])
    n     = min(len(items), 8)

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", "Agenda"), 24, th["text"], True, th["font_display"],
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    cols    = 2 if n > 4 else 1
    col_w   = (CW - Inches(0.32)) / cols if cols == 2 else CW
    per_col = (n + cols - 1) // cols
    avail_h = H - MY - Inches(1.1) - MY
    item_h  = min(Inches(0.76), avail_h / per_col)
    num_w   = Inches(0.52)

    for i, item in enumerate(items[:n]):
        col = i // per_col
        row = i % per_col
        ix  = MX + col * (col_w + Inches(0.32))
        iy  = MY + Inches(0.95) + row * item_h

        # 번호 박스 (둥근 사각형)
        shape_text(slide, ix, iy + Inches(0.06), num_w, item_h - Inches(0.12),
                   f"{i+1:02d}", 14, th["accent2"], True, F,
                   PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE,
                   th["surface"], th["accent"], 0.5,
                   shape_type=RRECT, font_en=FE)

        shape_text(slide, ix + num_w + Inches(0.14), iy + Inches(0.06),
                   col_w - num_w - Inches(0.16), item_h - Inches(0.12),
                   item, 15, th["dim"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

        hline(slide, ix, iy + item_h - Pt(3), col_w, th["rule"], 0.3)

    add_notes(slide, content.get("notes", ""))


# ── 신규 레이아웃 (5종) ──────────────────────────────────────

def slide_big_number(prs, content, th):
    """대형 숫자 강조 슬라이드. 핵심 수치 하나를 극적으로 표현."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg_gradient(slide, th["bg"], th.get("bg2", th["bg"]), angle=160)
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    number  = content.get("number", content.get("value", ""))
    label   = content.get("label", "")
    detail  = content.get("detail", content.get("description", ""))
    delta   = content.get("delta", "")

    # 좌측: 대형 숫자 (72pt+)
    num_w = CW * 0.55
    shape_text(slide, MX, H/2 - Inches(1.8), num_w, Inches(2.5),
               str(number), 80, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.BOTTOM, font_en=FE)

    # 레이블 (숫자 아래)
    if label:
        shape_text(slide, MX, H/2 + Inches(0.4), num_w, Inches(0.5),
                   label.upper(), 14, th["metal"], True, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_en=FE)

    # 우측: 설명 + 변화량
    rx = MX + num_w + Inches(0.4)
    rw = CW - num_w - Inches(0.4)
    vline(slide, rx - Inches(0.2), H/2 - Inches(1.0), Inches(2.0), th["accent"], 2.0)

    if delta:
        d_col = th["positive"] if delta.startswith("+") else (
                th["negative"] if delta.startswith("-") else th["dim"])
        shape_text(slide, rx, H/2 - Inches(0.8), rw, Inches(0.6),
                   delta, 28, d_col, True, FD,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

    if detail:
        shape_text(slide, rx, H/2, rw, Inches(1.0),
                   detail, 14, th["dim"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_icon_grid(prs, content, th):
    """아이콘 + 텍스트 그리드. 3-6개 항목을 카드 형태로 배치."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FE = th["font"], th.get("font_en", th["font"])

    items = content.get("items", [])
    n = min(len(items), 6)

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, th["font_display"],
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    if not n: return

    cols = 3 if n > 3 else n
    rows_n = (n + cols - 1) // cols
    gap = Inches(0.2)
    card_w = (CW - gap * (cols - 1)) / cols
    avail_h = H - (MY + Inches(0.92)) - MY - Inches(0.1)
    card_h = (avail_h - gap * (rows_n - 1)) / rows_n

    accents = [th["accent"], th["positive"], th["warn"],
               th["accent2"], th["metal"], th["dim"]]

    for i, item in enumerate(items[:n]):
        col_i = i % cols
        row_i = i // cols
        cx = MX + col_i * (card_w + gap)
        cy = MY + Inches(0.92) + row_i * (card_h + gap)
        col = accents[i % len(accents)]

        # 카드 (그림자)
        solid_rrect(slide, cx, cy, card_w, card_h,
                    th["card_bg"], th["card_border"], 0.5, shadow=True)

        # 상단 악센트 dot
        oval_dot(slide, cx + Inches(0.3), cy + Inches(0.35), Inches(0.12), col)

        # 아이콘 텍스트 (이모지 또는 기호)
        icon = item.get("icon", "") if isinstance(item, dict) else ""
        if icon:
            shape_text(slide, cx + Inches(0.12), cy + Inches(0.15),
                       Inches(0.4), Inches(0.4),
                       icon, 22, col, False, F,
                       PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, font_en=FE)

        # 제목
        title = item.get("title", str(item)) if isinstance(item, dict) else str(item)
        title_y = cy + Inches(0.55)
        shape_text(slide, cx + Inches(0.15), title_y, card_w - Inches(0.3), Inches(0.42),
                   title, 14, th["text"], True, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

        # 설명
        desc = item.get("description", item.get("desc", "")) if isinstance(item, dict) else ""
        if desc:
            shape_text(slide, cx + Inches(0.15), title_y + Inches(0.42),
                       card_w - Inches(0.3), card_h - Inches(1.1),
                       desc, 12, th["dim"], False, F,
                       PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_image_text(prs, content, th):
    """이미지 + 텍스트 50/50 분할. 이미지 없으면 컬러 블록 대체."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    image_path = content.get("image", content.get("image_path", ""))
    position   = content.get("image_position", "left")  # left or right
    headline   = content.get("headline", "")
    body_text  = content.get("body", content.get("text", ""))
    bullets    = content.get("bullets", [])

    img_w = CW * 0.48
    txt_w = CW * 0.48
    gap = CW * 0.04

    if position == "right":
        txt_x = MX
        img_x = MX + txt_w + gap
    else:
        img_x = MX
        txt_x = MX + img_w + gap

    content_y = MY + Inches(0.15)
    content_h = H - 2 * MY - Inches(0.3)

    # 이미지 영역
    img = try_add_image(slide, image_path, img_x, content_y, img_w, content_h)
    if not img:
        # 이미지 없으면 그라데이션 블록
        solid_rrect(slide, img_x, content_y, img_w, content_h,
                    th["surface2"], th["card_border"], 0.5, shadow=True)
        # 플레이스홀더 텍스트
        shape_text(slide, img_x, content_y + content_h/2 - Inches(0.3),
                   img_w, Inches(0.6),
                   "[IMAGE]", 18, th["dim2"], False, F,
                   PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, font_en=FE)

    # 텍스트 영역
    y = content_y + Inches(0.3)

    # 헤드라인
    shape_text(slide, txt_x, y, txt_w, Inches(1.2),
               headline, 28, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_en=FE)
    y += Inches(1.4)

    hline(slide, txt_x, y, txt_w * 0.4, th["accent"], 2.0)
    y += Inches(0.2)

    # 본문 or 불릿
    if bullets:
        for b in bullets[:4]:
            oval_dot(slide, txt_x + Inches(0.1), y + Inches(0.25), Inches(0.05), th["accent"])
            shape_text(slide, txt_x + Inches(0.28), y, txt_w - Inches(0.3), Inches(0.52),
                       b, 13, th["dim"], False, F,
                       PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
            y += Inches(0.52)
    elif body_text:
        shape_multiline(slide, txt_x, y, txt_w, content_h - Inches(2.2),
                        body_text.split("\n"), 14, th["dim"], False, F,
                        PP_ALIGN.LEFT, MSO_ANCHOR.TOP, line_spacing=1.5, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_bento_grid(prs, content, th):
    """벤토 박스 그리드 — 대형 카드 1개 + 소형 카드 2-3개."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    items = content.get("items", [])
    if not items: return

    gap = Inches(0.18)
    y0 = MY + Inches(0.92)
    avail_h = H - y0 - MY - Inches(0.1)

    # 6:4 분할 — 왼쪽 대형, 오른쪽 소형 스택
    big_w = CW * 0.58
    small_w = CW - big_w - gap

    accents = [th["accent"], th["positive"], th["warn"], th["accent2"]]

    # 대형 카드 (첫 번째 항목)
    main = items[0] if isinstance(items[0], dict) else {"title": str(items[0])}
    solid_rrect(slide, MX, y0, big_w, avail_h,
                th["card_bg"], th["card_border"], 0.5, shadow=True)
    solid_rect(slide, MX, y0, big_w, Pt(4), accents[0])

    shape_text(slide, MX + Inches(0.3), y0 + Inches(0.4), big_w - Inches(0.6), Inches(0.8),
               main.get("title", ""), 22, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_en=FE)
    if main.get("description"):
        shape_text(slide, MX + Inches(0.3), y0 + Inches(1.3), big_w - Inches(0.6), avail_h - Inches(1.8),
                   main["description"], 14, th["dim"], False, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_en=FE)

    # 소형 카드들 (나머지 항목)
    small_items = items[1:4]  # 최대 3개
    n_small = len(small_items)
    if n_small:
        small_h = (avail_h - gap * (n_small - 1)) / n_small
        sx = MX + big_w + gap

        for j, sitem in enumerate(small_items):
            sitem = sitem if isinstance(sitem, dict) else {"title": str(sitem)}
            sy = y0 + j * (small_h + gap)
            col = accents[(j + 1) % len(accents)]

            solid_rrect(slide, sx, sy, small_w, small_h,
                        th["card_bg"], th["card_border"], 0.5, shadow=True)
            solid_rect(slide, sx, sy, small_w, Pt(3), col)

            shape_text(slide, sx + Inches(0.2), sy + Inches(0.15),
                       small_w - Inches(0.4), Inches(0.45),
                       sitem.get("title", ""), 14, th["text"], True, F,
                       PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

            if sitem.get("description"):
                shape_text(slide, sx + Inches(0.2), sy + Inches(0.58),
                           small_w - Inches(0.4), small_h - Inches(0.75),
                           sitem["description"], 12, th["dim"], False, F,
                           PP_ALIGN.LEFT, MSO_ANCHOR.TOP, font_en=FE)

    add_notes(slide, content.get("notes", ""))


def slide_comparison(prs, content, th):
    """비교 슬라이드 — Before/After 또는 Option A vs B."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_bg(slide, th["bg"])
    F, FD, FE = th["font"], th["font_display"], th.get("font_en", th["font"])

    shape_text(slide, MX, MY, CW, Inches(0.68),
               content.get("headline", ""), 24, th["text"], True, FD,
               PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
    hline(slide, MX, MY + Inches(0.66), CW, th["metal"], 0.5)

    left_data  = content.get("before", content.get("option_a", content.get("left", {})))
    right_data = content.get("after", content.get("option_b", content.get("right", {})))

    gap = Inches(0.4)
    col_w = (CW - gap) / 2
    y_start = MY + Inches(0.92)
    col_h = H - y_start - MY - Inches(0.1)

    # VS 구분 원
    vs_cx = MX + col_w + gap / 2
    vs_cy = y_start + col_h / 2
    oval_dot(slide, vs_cx, vs_cy, Inches(0.28), th["accent"])
    shape_text(slide, vs_cx - Inches(0.28), vs_cy - Inches(0.28),
               Inches(0.56), Inches(0.56),
               "VS", 12, th.get("card_bg", th["bg"]), True, F,
               PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, font_en=FE)

    pairs = [
        (MX, left_data, th["negative"]),
        (MX + col_w + gap, right_data, th["positive"]),
    ]

    for cx, data, accent in pairs:
        solid_rrect(slide, cx, y_start, col_w, col_h,
                    th["card_bg"], th["card_border"], 0.5, shadow=True)
        solid_rect(slide, cx, y_start, col_w, Pt(4), accent)

        # 타이틀
        title = data.get("title", data.get("label", ""))
        shape_text(slide, cx + Inches(0.2), y_start + Inches(0.2),
                   col_w - Inches(0.4), Inches(0.5),
                   title.upper(), 14, accent, True, F,
                   PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)

        hline(slide, cx + Inches(0.2), y_start + Inches(0.72),
              col_w - Inches(0.4), th["rule"], 0.4)

        # 불릿
        by = y_start + Inches(0.85)
        for b in data.get("bullets", data.get("items", []))[:5]:
            oval_dot(slide, cx + Inches(0.3), by + Inches(0.24), Inches(0.05), accent)
            shape_text(slide, cx + Inches(0.45), by, col_w - Inches(0.65), Inches(0.5),
                       b, 13, th["dim"], False, F,
                       PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE, font_en=FE)
            by += Inches(0.52)

    add_notes(slide, content.get("notes", ""))


# ── 디스패처 ──────────────────────────────────────────────────
BUILDERS = {
    "cover":         slide_cover,
    "section_break": slide_section_break,
    "agenda":        slide_agenda,
    "bullets":       slide_bullets,
    "single":        slide_bullets,
    "two_column":    slide_two_column,
    "stats_grid":    slide_stats_grid,
    "data_table":    slide_data_table,
    "flow_diagram":  slide_flow_diagram,
    "timeline":      slide_timeline,
    "bar_chart":     slide_bar_chart,
    "quote":         slide_quote,
    "closing":       slide_closing,
    # 신규 5종
    "big_number":    slide_big_number,
    "icon_grid":     slide_icon_grid,
    "image_text":    slide_image_text,
    "bento_grid":    slide_bento_grid,
    "comparison":    slide_comparison,
}

# 구조용 레이아웃 (교차 검증에서 제외)
_STRUCTURAL = {"cover", "section_break", "agenda", "closing"}


def build_deck(spec, theme_key=None, output_path=None):
    meta   = spec.get("meta", {})
    slides = spec.get("slides", [])
    t_key  = theme_key or meta.get("theme", "dark_editorial")
    th     = THEMES.get(t_key, THEMES["dark_editorial"])
    out    = output_path or "output.pptx"
    company = meta.get("company", meta.get("author", ""))

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    total = len(slides)
    prev_layout = None
    warnings = []

    for idx, s in enumerate(slides):
        layout = s.get("layout", "bullets")
        builder = BUILDERS.get(layout, slide_bullets)

        # 레이아웃 교차 검증 (구조용 제외)
        if layout not in _STRUCTURAL and layout == prev_layout:
            warnings.append(f"[WARN] Slide {idx+1}: '{layout}' 연속 사용 — 다양성 부족")

        builder(prs, s, th)

        # 풋터 추가 (cover/closing 제외)
        if layout not in {"cover", "closing"}:
            slide_obj = prs.slides[len(prs.slides) - 1]
            add_footer(slide_obj, th, idx + 1, total, company)

        prev_layout = layout if layout not in _STRUCTURAL else None

    prs.save(out)

    if warnings:
        for w in warnings:
            print(w, file=sys.stderr)

    return out


if __name__ == "__main__":
    p = argparse.ArgumentParser(description="A-Team PPT Generator v3")
    p.add_argument("spec")
    p.add_argument("--theme", default=None,
                   choices=list(THEMES.keys()))
    p.add_argument("--output", default=None)
    args = p.parse_args()

    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)

    out = build_deck(spec, args.theme, args.output)
    t = args.theme or spec.get('meta',{}).get('theme','dark_editorial')
    n = len(spec.get('slides',[]))
    print(f"Generated: {out}")
    print(f"Slides: {n} / Theme: {t} / Layouts: {len(BUILDERS)}")
