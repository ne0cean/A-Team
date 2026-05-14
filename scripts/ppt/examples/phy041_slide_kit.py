"""
slide_kit — reusable slide type library for .pptx generation.
Content-agnostic. Pass in content dict + theme dict + variant, get a slide.

SLIDE TYPES (exported):
  structural:  cover, section_header, closing, agenda
  pitch:       problem, solution_mockup, why_now, bento_features,
               moat_columns, matrix_2x2, bar_chart, stats_grid,
               timeline, team, ask, quote
  academic:    motivation, research_question, method_diagram,
               results_chart, ablation_table, conclusions, references

THEMES: vc_clean, cyberpunk, editorial, academic_minimal, research_dark
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
CW = SLIDE_W - 2 * MARGIN


# ── Themes ────────────────────────────────────────────────────
THEMES = {
    "vc_clean": {
        "bg":          RGBColor(0xFA, 0xFA, 0xFA),
        "bg_alt":      RGBColor(0xF0, 0xF2, 0xF7),
        "text":        RGBColor(0x0A, 0x0A, 0x0A),
        "dim":         RGBColor(0x55, 0x55, 0x55),
        "dim2":        RGBColor(0x88, 0x88, 0x88),
        "accent":      RGBColor(0x25, 0x63, 0xEB),
        "accent2":     RGBColor(0x7C, 0x3A, 0xED),
        "positive":    RGBColor(0x00, 0x9E, 0x73),
        "card_bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "card_border": RGBColor(0xE1, 0xE4, 0xEA),
        "watermark":   RGBColor(0xEA, 0xED, 0xF3),
        "font":        "Helvetica",
        "font_display":"Helvetica",
    },
    "academic_minimal": {
        "bg":          RGBColor(0xFF, 0xFF, 0xFF),
        "bg_alt":      RGBColor(0xF5, 0xF5, 0xF5),
        "text":        RGBColor(0x10, 0x10, 0x10),
        "dim":         RGBColor(0x44, 0x44, 0x44),
        "dim2":        RGBColor(0x88, 0x88, 0x88),
        "accent":      RGBColor(0x1F, 0x4E, 0x79),
        "accent2":     RGBColor(0x2E, 0x75, 0xB6),
        "positive":    RGBColor(0x2F, 0x74, 0x2F),
        "card_bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "card_border": RGBColor(0xDA, 0xDC, 0xDF),
        "watermark":   RGBColor(0xF0, 0xF2, 0xF5),
        "font":        "Helvetica",
        "font_display":"Georgia",
    },
    "research_dark": {
        "bg":          RGBColor(0x0A, 0x0F, 0x1E),
        "bg_alt":      RGBColor(0x14, 0x1A, 0x2E),
        "text":        RGBColor(0xF0, 0xF4, 0xFA),
        "dim":         RGBColor(0xA8, 0xB2, 0xC8),
        "dim2":        RGBColor(0x6C, 0x76, 0x8C),
        "accent":      RGBColor(0x5C, 0xE1, 0xE6),
        "accent2":     RGBColor(0xFF, 0x8C, 0x42),
        "positive":    RGBColor(0x00, 0xE5, 0xA0),
        "card_bg":     RGBColor(0x10, 0x18, 0x2E),
        "card_border": RGBColor(0x2A, 0x35, 0x50),
        "watermark":   RGBColor(0x14, 0x1E, 0x38),
        "font":        "Helvetica",
        "font_display":"Helvetica",
    },
    "editorial": {
        "bg":          RGBColor(0xFA, 0xF5, 0xEA),
        "bg_alt":      RGBColor(0xF0, 0xE7, 0xD4),
        "text":        RGBColor(0x1A, 0x0F, 0x08),
        "dim":         RGBColor(0x6B, 0x4F, 0x3C),
        "dim2":        RGBColor(0x9A, 0x80, 0x6C),
        "accent":      RGBColor(0xB2, 0x35, 0x1E),
        "accent2":     RGBColor(0x5C, 0x3A, 0x1E),
        "positive":    RGBColor(0x4A, 0x6B, 0x2E),
        "card_bg":     RGBColor(0xFF, 0xFC, 0xF5),
        "card_border": RGBColor(0xD9, 0xC9, 0xAE),
        "watermark":   RGBColor(0xF0, 0xE5, 0xCD),
        "font":        "Georgia",
        "font_display":"Georgia",
    },
}


# ── Primitives ────────────────────────────────────────────────
def _set_font(run, size, color=None, bold=False, name="Helvetica", italic=False):
    run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name


def paint_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def tbox(slide, l, t, w, h, text, size=14, color=None, bold=False,
         name="Helvetica", alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         italic=False, line_spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = alignment
    if line_spacing is not None:
        p.line_spacing = line_spacing
    r = p.add_run(); r.text = text
    _set_font(r, size, color, bold, name, italic)
    return tb


def multi_para(slide, l, t, w, h, paragraphs, size=40, bold=True,
               name="Helvetica", alignment=PP_ALIGN.LEFT, line_spacing=1.1,
               anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, segs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment; p.line_spacing = line_spacing
        for seg in segs:
            text, color = seg[0], seg[1]
            opts = seg[2] if len(seg) > 2 else {}
            r = p.add_run(); r.text = text
            _set_font(r, size, color, opts.get("bold", bold), name,
                      opts.get("italic", False))
    return tb


def rect(slide, l, t, w, h, fill, border=None, bw=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw or 0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def rrect(slide, l, t, w, h, fill, border=None, bw=0, radius=0.04):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bw or 0.75)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    s.adjustments[0] = radius
    return s


def oval(slide, cx, cy, r, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, 2*r, 2*r)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def hline(slide, l, t, w, color, weight=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(weight))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def vline(slide, l, t, h, color, weight=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(weight), h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def corner_accent(slide, theme, variant="clean"):
    if variant == "neon":
        hline(slide, Inches(0.7), Inches(0.35), CW, theme["accent"], 1.5)
        hline(slide, Inches(0.7), Inches(7.15), CW, theme["accent2"], 1.5)
    elif variant == "rule":
        hline(slide, Inches(0.7), Inches(0.55), CW, theme["dim2"], 0.5)
        hline(slide, Inches(0.7), Inches(7.05), CW, theme["dim2"], 0.5)
    else:
        rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, theme["accent"])


def kicker(slide, theme, text, top=Inches(0.55)):
    tbox(slide, MARGIN, top, CW, Inches(0.28), text,
         size=11, color=theme["accent"], bold=True, name=theme["font"])


def footer(slide, theme, current, total, author="", accent_variant="clean"):
    if accent_variant != "neon":
        pass
    tbox(slide, MARGIN, Inches(7.15), Inches(8), Inches(0.28),
         author, size=9, color=theme["dim2"], name=theme["font"])
    tbox(slide, Inches(11.2), Inches(7.15), Inches(1.4), Inches(0.28),
         f"{current} / {total}", size=9, color=theme["dim2"],
         name=theme["font"], alignment=PP_ALIGN.RIGHT)


def chip(slide, l, t, text, fill, text_color, name="Helvetica"):
    w = Inches(len(text) * 0.09 + 0.4); h = Inches(0.3)
    rrect(slide, l, t, w, h, fill, radius=0.35)
    tbox(slide, l, t + Inches(0.02), w, h, text, size=10,
         color=text_color, bold=True, name=name, alignment=PP_ALIGN.CENTER)


# ── Slide builders ────────────────────────────────────────────
def build_cover(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, h1_line1, h1_line2_accent, lede, chips?: [(text,fill_key,text_key)]}"""
    paint_bg(slide, theme["bg"])
    if variant == "clean":
        oval(slide, Inches(12.3), Inches(1.5), Inches(0.8), theme["accent2"])
    elif variant == "neon":
        hline(slide, Inches(0), Inches(3.7), SLIDE_W, theme["accent"], 0.3)
        vline(slide, Inches(6.6), Inches(0), SLIDE_H, theme["accent2"], 0.3)
        oval(slide, Inches(6.6), Inches(3.7), Inches(0.12), theme["bg"], theme["accent"])
    elif variant == "rule":
        hline(slide, Inches(5.5), Inches(2.5), Inches(2.3), theme["accent"], 1.0)
    # Optional logomark
    if c.get("logo"):
        lx, ly = Inches(6.3), Inches(1.35)
        rrect(slide, lx, ly, Inches(0.7), Inches(0.7), theme["accent"], radius=0.2)
        tbox(slide, lx, ly, Inches(0.7), Inches(0.7), c["logo"],
             size=36, color=theme["bg"], bold=True, name=theme["font_display"],
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Kicker
    tbox(slide, MARGIN, Inches(2.55), CW, Inches(0.3), c["kicker"],
         size=12, color=theme["accent"], bold=True, name=theme["font"],
         alignment=PP_ALIGN.CENTER)
    # Headline
    multi_para(slide, MARGIN, Inches(3.05), CW, Inches(2.5),
               [[(c["h1_line1"], theme["text"])],
                [(c["h1_line2_accent"], theme["accent"])]],
               size=54, bold=True, name=theme["font_display"],
               alignment=PP_ALIGN.CENTER, line_spacing=1.15)
    # Lede
    tbox(slide, MARGIN, Inches(5.95), CW, Inches(0.4), c["lede"],
         size=16, color=theme["dim"], name=theme["font"],
         alignment=PP_ALIGN.CENTER)
    # Chips (resolve theme-key strings to RGBColor)
    def _resolve(val):
        if isinstance(val, str):
            if val in theme:
                return theme[val]
            if len(val) == 6 and all(ch in "0123456789abcdefABCDEF" for ch in val):
                return RGBColor(int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
        return val
    if c.get("chips"):
        chips = c["chips"]
        cw = [Inches(len(t[0]) * 0.09 + 0.4) for t in chips]
        gap = Inches(0.2)
        tw = sum(cw, Inches(0)) + gap * (len(chips) - 1)
        cur = MARGIN + (CW - tw) / 2
        for (text, fill, col), w in zip(chips, cw):
            chip(slide, cur, Inches(6.5), text,
                 _resolve(fill), _resolve(col), theme["font"])
            cur += w + gap


def build_section_header(slide, c, theme, variant="clean", ctx=None):
    """c: {num, title, subtitle?}"""
    paint_bg(slide, theme["bg"])
    # Giant numeral — size dialed to fit 2 chars in wide box
    tbox(slide, Inches(-0.3), Inches(1.0), Inches(6.5), Inches(5.5), c["num"],
         size=300, color=theme["watermark"], bold=True,
         name=theme["font_display"], alignment=PP_ALIGN.CENTER)
    hline(slide, Inches(6.2), Inches(2.85), Inches(1.5), theme["accent"], 2.5)
    tbox(slide, Inches(6.2), Inches(3.0), Inches(7.0), Inches(1.5), c["title"],
         size=50, color=theme["text"], bold=True, name=theme["font_display"])
    if c.get("subtitle"):
        tbox(slide, Inches(6.2), Inches(4.3), Inches(7.0), Inches(1.0),
             c["subtitle"], size=18, color=theme["dim"], name=theme["font"],
             italic=True)


def build_problem(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title_line1_segs, title_line2, lede, cards:[(big,small)]}"""
    paint_bg(slide, theme["bg"])
    corner_accent(slide, theme, variant)
    tbox(slide, Inches(7.0), Inches(-0.5), Inches(6.5), Inches(5), c.get("watermark", "01"),
         size=380, color=theme["watermark"], bold=True,
         name=theme["font_display"], alignment=PP_ALIGN.RIGHT)
    kicker(slide, theme, c["kicker"])
    # Title — support either segs or plain
    title_segs = c.get("title_line1_segs") or [(c.get("title_line1", ""), theme["text"])]
    title_line2 = c.get("title_line2", "")
    paragraphs = [[(t, theme.get(k, k) if isinstance(k, str) else k) for t, k in title_segs]]
    if title_line2:
        paragraphs.append([(title_line2, theme["text"])])
    multi_para(slide, MARGIN, Inches(1.0), Inches(8), Inches(2.2),
               paragraphs, size=42, bold=True, name=theme["font_display"],
               line_spacing=1.1)
    tbox(slide, MARGIN, Inches(3.4), Inches(8.3), Inches(1.0), c["lede"],
         size=14, color=theme["dim"], name=theme["font"], line_spacing=1.4)
    # Cards
    cards = c["cards"]; n = len(cards)
    card_w = (CW - Inches(0.15 * (n - 1))) / n
    y = Inches(4.7); h = Inches(1.9)
    for i, (big, small) in enumerate(cards):
        x = MARGIN + i * (card_w + Inches(0.15))
        rrect(slide, x, y, card_w, h, theme["card_bg"], theme["card_border"])
        rect(slide, x, y, Inches(0.12), h, theme["accent"])
        tbox(slide, x + Inches(0.35), y + Inches(0.25),
             card_w - Inches(0.5), Inches(0.9), big,
             size=30, color=theme["accent"], bold=True, name=theme["font_display"])
        tbox(slide, x + Inches(0.35), y + Inches(1.1),
             card_w - Inches(0.5), Inches(0.7), small,
             size=12, color=theme["dim"], name=theme["font"])


def build_why_now(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, lede, forces:[(badge,title,body)]}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    tbox(slide, Inches(-0.5), Inches(-0.5), Inches(6), Inches(5),
         c.get("watermark", "02"), size=380, color=theme["watermark"],
         bold=True, name=theme["font_display"], alignment=PP_ALIGN.LEFT)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(1.2), c["title"],
         size=40, color=theme["text"], bold=True, name=theme["font_display"])
    tbox(slide, MARGIN, Inches(2.2), CW, Inches(0.6), c["lede"],
         size=14, color=theme["dim"], name=theme["font"])
    col_w = (CW - Inches(0.6)) / 3; y = Inches(3.5)
    for i, (badge, title, body) in enumerate(c["forces"]):
        x = MARGIN + i * (col_w + Inches(0.3))
        rrect(slide, x, y, Inches(0.9), Inches(0.35), theme["accent"], radius=0.3)
        tbox(slide, x, y + Inches(0.03), Inches(0.9), Inches(0.3), badge,
             size=10, color=theme["bg"], bold=True, name=theme["font"],
             alignment=PP_ALIGN.CENTER)
        tbox(slide, x, y + Inches(0.65), col_w, Inches(0.6), title,
             size=20, color=theme["text"], bold=True, name=theme["font_display"])
        tbox(slide, x, y + Inches(1.4), col_w, Inches(2.0), body,
             size=12, color=theme["dim"], name=theme["font"], line_spacing=1.3)
        hline(slide, x, y + Inches(3.5), col_w, theme["accent2"], 1.0)


def build_bento_features(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, flagship:{badge,title,body}, sub:[(badge,title,body)x3]}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=34, color=theme["text"], bold=True, name=theme["font_display"])
    y = Inches(2.4)
    # Flagship
    rrect(slide, MARGIN, y, Inches(6.0), Inches(4.3),
          theme["card_bg"], theme["card_border"])
    rect(slide, MARGIN, y, Inches(0.12), Inches(4.3), theme["accent"])
    f = c["flagship"]
    tbox(slide, MARGIN + Inches(0.4), y + Inches(0.4), Inches(5.3), Inches(0.4),
         f["badge"], size=10, color=theme["accent"], bold=True, name=theme["font"])
    tbox(slide, MARGIN + Inches(0.4), y + Inches(0.8), Inches(5.3), Inches(0.7),
         f["title"], size=28, color=theme["text"], bold=True,
         name=theme["font_display"])
    tbox(slide, MARGIN + Inches(0.4), y + Inches(1.6), Inches(5.3), Inches(2.5),
         f["body"], size=13, color=theme["dim"], name=theme["font"], line_spacing=1.4)
    # Sub cards
    cx = Inches(7.0); cw = CW - Inches(6.3)
    heights = [Inches(1.3), Inches(1.3), Inches(1.5)]
    yy = y
    for (badge, title, body), ch in zip(c["sub"], heights):
        rrect(slide, cx, yy, cw, ch, theme["card_bg"], theme["card_border"])
        tbox(slide, cx + Inches(0.3), yy + Inches(0.2), Inches(2), Inches(0.3),
             badge, size=9, color=theme["accent2"], bold=True, name=theme["font"])
        tbox(slide, cx + Inches(0.3), yy + Inches(0.5), cw - Inches(0.6), Inches(0.5),
             title, size=19, color=theme["text"], bold=True, name=theme["font_display"])
        tbox(slide, cx + Inches(0.3), yy + Inches(1.05), cw - Inches(0.6), Inches(0.6),
             body, size=11, color=theme["dim"], name=theme["font"])
        yy += ch + Inches(0.1)


def build_moat_columns(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, cards:[(num,title,body)]}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=36, color=theme["text"], bold=True, name=theme["font_display"])
    n = len(c["cards"]); card_w = (CW - Inches(0.2 * (n - 1))) / n
    y = Inches(2.6); h = Inches(3.8)
    for i, (num, title, body) in enumerate(c["cards"]):
        x = MARGIN + i * (card_w + Inches(0.2))
        rrect(slide, x, y, card_w, h, theme["card_bg"], theme["card_border"])
        tbox(slide, x, y + Inches(0.2), card_w, Inches(1.5), num,
             size=100, color=theme["watermark"], bold=True,
             name=theme["font_display"], alignment=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.3), y + Inches(1.8),
             card_w - Inches(0.6), Inches(0.6), title,
             size=20, color=theme["text"], bold=True, name=theme["font_display"])
        tbox(slide, x + Inches(0.3), y + Inches(2.5),
             card_w - Inches(0.6), Inches(1.2), body,
             size=12, color=theme["dim"], name=theme["font"], line_spacing=1.4)


def build_matrix_2x2(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, x_axis:(low,high), y_axis:(low,high),
          competitors:[{name, x:0-1, y:0-1, is_us:bool}]}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=30, color=theme["text"], bold=True, name=theme["font_display"])
    mx, my = Inches(3.0), Inches(2.3); mw, mh = Inches(7.5), Inches(4.3)
    rrect(slide, mx, my, mw, mh, theme["card_bg"], theme["card_border"])
    hline(slide, mx, my + mh / 2, mw, theme["dim2"], 0.5)
    vline(slide, mx + mw / 2, my, mh, theme["dim2"], 0.5)
    # axis labels
    tbox(slide, mx, my - Inches(0.45), mw, Inches(0.3),
         f"{c['y_axis'][1]} ▲", size=11, color=theme["dim"],
         name=theme["font"], alignment=PP_ALIGN.CENTER, italic=True)
    tbox(slide, mx, my + mh + Inches(0.1), mw, Inches(0.3),
         f"{c['y_axis'][0]} ▼", size=11, color=theme["dim"],
         name=theme["font"], alignment=PP_ALIGN.CENTER, italic=True)
    tbox(slide, Inches(0.3), my + mh / 2 - Inches(0.2),
         Inches(2.5), Inches(0.4), f"{c['x_axis'][0]} ◀",
         size=11, color=theme["dim"], name=theme["font"],
         alignment=PP_ALIGN.RIGHT, italic=True)
    tbox(slide, mx + mw + Inches(0.2), my + mh / 2 - Inches(0.2),
         Inches(2.5), Inches(0.4), f"▶ {c['x_axis'][1]}",
         size=11, color=theme["dim"], name=theme["font"], italic=True)
    for comp in c["competitors"]:
        cx = mx + mw * comp["x"]; cy = my + mh * (1 - comp["y"])
        r = Inches(0.32 if comp.get("is_us") else 0.20)
        fill = theme["accent"] if comp.get("is_us") else theme["dim"]
        oval(slide, cx, cy, r, fill)
        lbl = f"★ {comp['name']}" if comp.get("is_us") else comp["name"]
        size = 16 if comp.get("is_us") else 12
        color = theme["accent"] if comp.get("is_us") else theme["text"]
        tbox(slide, cx - Inches(1.2), cy + r + Inches(0.05),
             Inches(2.4), Inches(0.3), lbl, size=size,
             color=color, bold=True,
             name=theme["font_display"] if comp.get("is_us") else theme["font"],
             alignment=PP_ALIGN.CENTER)


def build_bar_chart(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, subtitle, categories:[str], values:[num], series_name, source}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(1.0), c["title"],
         size=34, color=theme["text"], bold=True, name=theme["font_display"])
    tbox(slide, MARGIN, Inches(2.0), CW, Inches(0.5), c.get("subtitle", ""),
         size=13, color=theme["dim"], name=theme["font"], italic=True)
    chart_data = CategoryChartData()
    chart_data.categories = c["categories"]
    chart_data.add_series(c.get("series_name", "Value"), c["values"])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, Inches(2.6),
        CW, Inches(3.8), chart_data
    ).chart
    chart.has_title = False; chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(14)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = theme["text"]
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = theme["accent"]
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(11)
    ca.tick_labels.font.color.rgb = theme["dim"]
    ca.format.line.color.rgb = theme["dim2"]
    va = chart.value_axis; va.visible = False
    tbox(slide, MARGIN, Inches(6.6), CW, Inches(0.4),
         c.get("source", ""), size=9, color=theme["dim2"],
         name=theme["font"], italic=True)


def build_stats_grid(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, stats:[(big, mid, note)]}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=36, color=theme["text"], bold=True, name=theme["font_display"])
    n = len(c["stats"]); w = (CW - Inches(0.2 * (n - 1))) / n
    y = Inches(2.6); h = Inches(3.5)
    for i, (big, mid, note) in enumerate(c["stats"]):
        x = MARGIN + i * (w + Inches(0.2))
        rrect(slide, x, y, w, h, theme["card_bg"], theme["card_border"])
        tbox(slide, x, y + Inches(0.6), w, Inches(1.6), big,
             size=52, color=theme["accent"], bold=True,
             name=theme["font_display"], alignment=PP_ALIGN.CENTER)
        tbox(slide, x, y + Inches(2.1), w, Inches(0.5), mid,
             size=16, color=theme["text"], bold=True,
             name=theme["font_display"], alignment=PP_ALIGN.CENTER)
        tbox(slide, x, y + Inches(2.7), w, Inches(0.5), note,
             size=11, color=theme["dim"], name=theme["font"],
             italic=True, alignment=PP_ALIGN.CENTER)


def build_timeline(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, milestones:[(label, caption)]}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=34, color=theme["text"], bold=True, name=theme["font_display"])
    n = len(c["milestones"])
    track_y = Inches(4.0)
    hline(slide, MARGIN + Inches(0.5), track_y, CW - Inches(1.0),
          theme["accent"], 2.5)
    step = (CW - Inches(1.0)) / (n - 1) if n > 1 else Inches(0)
    for i, (label, caption) in enumerate(c["milestones"]):
        cx = MARGIN + Inches(0.5) + step * i
        oval(slide, cx, track_y + Inches(0.015), Inches(0.15),
             theme["bg"], theme["accent"])
        oval(slide, cx, track_y + Inches(0.015), Inches(0.08),
             theme["accent"])
        # alternating above/below
        above = (i % 2 == 0)
        label_y = track_y - Inches(1.5) if above else track_y + Inches(0.3)
        caption_y = label_y + Inches(0.5)
        tbox(slide, cx - Inches(1.2), label_y, Inches(2.4), Inches(0.4),
             label, size=14, color=theme["text"], bold=True,
             name=theme["font_display"], alignment=PP_ALIGN.CENTER)
        tbox(slide, cx - Inches(1.2), caption_y, Inches(2.4), Inches(0.7),
             caption, size=10, color=theme["dim"], name=theme["font"],
             alignment=PP_ALIGN.CENTER)


def build_team(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, members:[(name, role, bio, initials)]}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=34, color=theme["text"], bold=True, name=theme["font_display"])
    n = len(c["members"]); w = (CW - Inches(0.3 * (n - 1))) / n
    y = Inches(2.6); h = Inches(4.0)
    for i, (member_name, role, bio, initials) in enumerate(c["members"]):
        x = MARGIN + i * (w + Inches(0.3))
        rrect(slide, x, y, w, h, theme["card_bg"], theme["card_border"])
        avatar_cx = x + w / 2; avatar_cy = y + Inches(1.0)
        oval(slide, avatar_cx, avatar_cy, Inches(0.7), theme["accent2"])
        tbox(slide, x, y + Inches(0.65), w, Inches(0.7), initials,
             size=28, color=theme["bg"], bold=True,
             name=theme["font_display"], alignment=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.2), y + Inches(2.0),
             w - Inches(0.4), Inches(0.5), member_name,
             size=18, color=theme["text"], bold=True,
             name=theme["font_display"], alignment=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.2), y + Inches(2.55),
             w - Inches(0.4), Inches(0.4), role,
             size=11, color=theme["accent"], bold=True, name=theme["font"],
             alignment=PP_ALIGN.CENTER)
        tbox(slide, x + Inches(0.2), y + Inches(3.05),
             w - Inches(0.4), Inches(1.0), bio,
             size=11, color=theme["dim"], name=theme["font"],
             alignment=PP_ALIGN.CENTER, italic=True, line_spacing=1.3)


def build_ask(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, amount, raise_title:[(prefix,text),(amount,accent),(suffix,text)],
          lede, uses:[(pct, head, body)], milestone, footer_note}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    # "Raising $X seed"
    multi_para(slide, MARGIN, Inches(1.0), CW, Inches(2.0),
               [[(c["title_prefix"], theme["text"]),
                 (c["title_accent"], theme["accent"]),
                 (c["title_suffix"], theme["text"])]],
               size=54, bold=True, name=theme["font_display"])
    y = Inches(3.2)
    n = len(c["uses"]); col_w = (CW - Inches(0.3 * (n - 1))) / n
    for i, (pct, head, body) in enumerate(c["uses"]):
        x = MARGIN + i * (col_w + Inches(0.3))
        tbox(slide, x, y, col_w, Inches(0.7), pct,
             size=36, color=theme["accent"], bold=True, name=theme["font_display"])
        tbox(slide, x, y + Inches(0.8), col_w, Inches(0.4), head,
             size=18, color=theme["text"], bold=True, name=theme["font_display"])
        bar_y = y + Inches(1.3)
        rect(slide, x, bar_y, col_w, Inches(0.08), theme["bg_alt"])
        try:
            pct_val = int(pct.strip("%"))
        except ValueError:
            pct_val = 50
        rect(slide, x, bar_y, col_w * pct_val / 100, Inches(0.08), theme["accent"])
        tbox(slide, x, y + Inches(1.5), col_w, Inches(1.2), body,
             size=12, color=theme["dim"], name=theme["font"], line_spacing=1.4)
    if c.get("milestone"):
        tbox(slide, MARGIN, Inches(6.0), CW, Inches(0.4), c["milestone"],
             size=13, color=theme["text"], bold=True, name=theme["font"])
    if c.get("footer_note"):
        tbox(slide, MARGIN, Inches(6.45), CW, Inches(0.4), c["footer_note"],
             size=11, color=theme["dim"], name=theme["font"], italic=True)


def build_closing(slide, c, theme, variant="clean", ctx=None):
    """c: {eyebrow, headline, sub, contact, attribution}"""
    paint_bg(slide, theme["bg"])
    tbox(slide, Inches(-0.5), Inches(-0.5), Inches(14), Inches(7),
         c.get("eyebrow", ""), size=22, color=theme["accent"], bold=True,
         name=theme["font_display"], alignment=PP_ALIGN.CENTER)
    tbox(slide, MARGIN, Inches(2.5), CW, Inches(2.0), c["headline"],
         size=56, color=theme["text"], bold=True,
         name=theme["font_display"], alignment=PP_ALIGN.CENTER)
    if c.get("sub"):
        tbox(slide, MARGIN, Inches(3.9), CW, Inches(0.8), c["sub"],
             size=24, color=theme["accent"], bold=True,
             name=theme["font_display"], alignment=PP_ALIGN.CENTER,
             italic=True)
    tbox(slide, MARGIN, Inches(5.5), CW, Inches(0.5), c["contact"],
         size=18, color=theme["text"], bold=True, name=theme["font"],
         alignment=PP_ALIGN.CENTER)
    if c.get("attribution"):
        tbox(slide, MARGIN, Inches(6.1), CW, Inches(0.4), c["attribution"],
             size=13, color=theme["dim"], name=theme["font"],
             alignment=PP_ALIGN.CENTER, italic=True)


def build_quote(slide, c, theme, variant="clean", ctx=None):
    """c: {quote, attribution}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    tbox(slide, Inches(0.5), Inches(0.8), Inches(2), Inches(3),
         "\u201C", size=280, color=theme["watermark"], bold=True,
         name=theme["font_display"], alignment=PP_ALIGN.LEFT)
    tbox(slide, Inches(2.5), Inches(2.5), Inches(10), Inches(2.5),
         c["quote"], size=32, color=theme["text"], bold=False,
         italic=True, name=theme["font_display"], line_spacing=1.3)
    hline(slide, Inches(2.5), Inches(5.2), Inches(0.8), theme["accent"], 2)
    tbox(slide, Inches(2.5), Inches(5.4), Inches(10), Inches(0.5),
         c["attribution"], size=14, color=theme["dim"], bold=True,
         name=theme["font"])


def build_method_diagram(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, image_path?, caption, steps?:[str]}
    If image_path given, embed it full-width. Otherwise render text-based steps.
    """
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=32, color=theme["text"], bold=True, name=theme["font_display"])
    img_top = Inches(2.2); img_h = Inches(4.0)
    if c.get("image_path"):
        import os
        if os.path.exists(c["image_path"]):
            slide.shapes.add_picture(c["image_path"], MARGIN, img_top,
                                     width=CW, height=img_h)
        else:
            rrect(slide, MARGIN, img_top, CW, img_h, theme["bg_alt"],
                  theme["card_border"])
            tbox(slide, MARGIN, img_top + Inches(1.5), CW, Inches(0.5),
                 f"[image missing: {c['image_path']}]",
                 size=14, color=theme["dim"], name=theme["font"],
                 alignment=PP_ALIGN.CENTER, italic=True)
    elif c.get("steps"):
        # Render boxes with arrows
        n = len(c["steps"])
        box_w = (CW - Inches(0.4 * (n - 1))) / n
        for i, step in enumerate(c["steps"]):
            x = MARGIN + i * (box_w + Inches(0.4))
            rrect(slide, x, img_top + Inches(1.3), box_w, Inches(1.4),
                  theme["card_bg"], theme["card_border"])
            tbox(slide, x + Inches(0.2), img_top + Inches(1.5),
                 box_w - Inches(0.4), Inches(1.0), step,
                 size=13, color=theme["text"], bold=True,
                 name=theme["font_display"], alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
            if i < n - 1:
                tbox(slide, x + box_w, img_top + Inches(1.8),
                     Inches(0.4), Inches(0.4), "→",
                     size=24, color=theme["accent"], bold=True,
                     name=theme["font"], alignment=PP_ALIGN.CENTER)
    if c.get("caption"):
        tbox(slide, MARGIN, Inches(6.4), CW, Inches(0.6), c["caption"],
             size=12, color=theme["dim"], name=theme["font"], italic=True,
             alignment=PP_ALIGN.CENTER)


def build_results_table(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, headers:[str], rows:[[str]], highlight_col?:int, takeaway?:str}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=30, color=theme["text"], bold=True, name=theme["font_display"])
    nrows = len(c["rows"]) + 1
    ncols = len(c["headers"])
    t_top = Inches(2.3); t_h = Inches(3.8)
    tbl_shape = slide.shapes.add_table(nrows, ncols, MARGIN, t_top, CW, t_h)
    tbl = tbl_shape.table
    # Header row
    for j, h in enumerate(c["headers"]):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["accent"]
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        _set_font(r, 14, theme["bg"], True, theme["font_display"])
    # Data rows
    hc = c.get("highlight_col")
    for i, row in enumerate(c["rows"]):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            is_highlight = (hc is not None and j == hc)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme["bg_alt"] if is_highlight else theme["card_bg"]
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            _set_font(r, 13,
                      theme["accent"] if is_highlight else theme["text"],
                      is_highlight, theme["font"])
    if c.get("takeaway"):
        tbox(slide, MARGIN, Inches(6.4), CW, Inches(0.5), c["takeaway"],
             size=13, color=theme["accent"], bold=True,
             name=theme["font"], italic=True,
             alignment=PP_ALIGN.CENTER)


def build_references(slide, c, theme, variant="clean", ctx=None):
    """c: {kicker, title, refs:[str]}"""
    paint_bg(slide, theme["bg"]); corner_accent(slide, theme, variant)
    kicker(slide, theme, c["kicker"])
    tbox(slide, MARGIN, Inches(1.0), CW, Inches(0.9), c["title"],
         size=30, color=theme["text"], bold=True, name=theme["font_display"])
    y = Inches(2.2)
    for i, ref in enumerate(c["refs"]):
        tbox(slide, MARGIN, y + Inches(0.4 * i), Inches(0.5), Inches(0.35),
             f"[{i+1}]", size=11, color=theme["accent"], bold=True, name=theme["font"])
        tbox(slide, MARGIN + Inches(0.5), y + Inches(0.4 * i),
             CW - Inches(0.5), Inches(0.35), ref,
             size=11, color=theme["text"], name=theme["font"])


# ── Dispatcher ────────────────────────────────────────────────
SLIDE_TYPES = {
    "cover": build_cover,
    "section_header": build_section_header,
    "problem": build_problem,
    "why_now": build_why_now,
    "bento_features": build_bento_features,
    "moat_columns": build_moat_columns,
    "matrix_2x2": build_matrix_2x2,
    "bar_chart": build_bar_chart,
    "stats_grid": build_stats_grid,
    "timeline": build_timeline,
    "team": build_team,
    "ask": build_ask,
    "closing": build_closing,
    "quote": build_quote,
    "method_diagram": build_method_diagram,
    "results_table": build_results_table,
    "references": build_references,
}


def build_deck(slide_defs, theme_key, output_path, variant="clean", author=""):
    """slide_defs: list of (slide_type, content_dict) tuples."""
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    theme = THEMES[theme_key]
    total = len(slide_defs)
    for idx, (stype, content) in enumerate(slide_defs, 1):
        slide = prs.slides.add_slide(blank)
        fn = SLIDE_TYPES[stype]
        fn(slide, content, theme, variant, ctx={"idx": idx, "total": total})
        # Footer except on cover / closing / section_header
        if stype not in ("cover", "closing", "section_header"):
            footer(slide, theme, idx, total, author=author, accent_variant=variant)
    prs.save(output_path)
    return output_path
