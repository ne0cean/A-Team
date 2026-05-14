#!/usr/bin/env python3
"""
PPT Generator — A-Team
Professional .pptx generator with Korean-native support.

Usage:
    python generate.py spec.json
    python generate.py spec.json --output deck.pptx
    python generate.py spec.json --theme consulting|minimal|executive

Requirements:
    pip install python-pptx
"""

import argparse
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Slide dimensions (16:9 widescreen) ─────────────────────────
SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)
MARGIN  = Cm(1.8)

# ─── Themes ──────────────────────────────────────────────────────
THEMES = {
    "consulting": {
        "bg":        (248, 250, 252),
        "primary":   ( 15,  23,  42),
        "accent":    ( 37,  99, 235),
        "secondary": (100, 116, 139),
        "rule":      (203, 213, 225),
        "surface":   (241, 245, 249),
        "white":     (255, 255, 255),
        "font":      "Malgun Gothic",
    },
    "minimal": {
        "bg":        (255, 255, 255),
        "primary":   ( 24,  24,  27),
        "accent":    ( 24,  24,  27),
        "secondary": (113, 113, 122),
        "rule":      (228, 228, 231),
        "surface":   (244, 244, 245),
        "white":     (255, 255, 255),
        "font":      "Malgun Gothic",
    },
    "executive": {
        "bg":        (250, 250, 249),
        "primary":   ( 28,  25,  23),
        "accent":    (159,  18,  57),
        "secondary": (120, 113, 108),
        "rule":      (214, 211, 209),
        "surface":   (245, 245, 244),
        "white":     (255, 255, 255),
        "font":      "Malgun Gothic",
    },
}


# ─── Helpers ──────────────────────────────────────────────────────

def rgb(t):
    return RGBColor(t[0], t[1], t[2])


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, font, size, color,
             bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_bullets(slide, x, y, w, h, bullets, font, size, primary, secondary, bullet="—"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        p.space_after = Pt(3)

        r1 = p.add_run()
        r1.text = bullet + "  "
        r1.font.name = font
        r1.font.size = Pt(size)
        r1.font.color.rgb = rgb(secondary)

        r2 = p.add_run()
        r2.text = item
        r2.font.name = font
        r2.font.size = Pt(size)
        r2.font.color.rgb = rgb(primary)

    return txBox


def add_slide_number(slide, num, t):
    if num <= 1:
        return
    add_text(slide,
        SLIDE_W - Cm(2.8), SLIDE_H - Cm(1.0), Cm(2.5), Cm(0.8),
        str(num), t["font"], 9, t["secondary"],
        align=PP_ALIGN.RIGHT
    )


def title_bar(slide, title, t):
    add_rect(slide, MARGIN, Cm(1.5), Cm(0.2), Cm(1.2), t["accent"])
    add_text(slide,
        MARGIN + Cm(0.6), Cm(1.3), SLIDE_W - MARGIN * 2 - Cm(0.6), Cm(1.6),
        title, t["font"], 24, t["primary"], bold=True
    )
    add_rect(slide, MARGIN, Cm(3.2), SLIDE_W - MARGIN * 2, Cm(0.04), t["rule"])


def add_speaker_notes(slide, text):
    if not text:
        return
    ns = slide.notes_slide
    ns.notes_text_frame.text = text


# ─── Layout builders ─────────────────────────────────────────────

def build_cover(slide, data, t, meta):
    set_bg(slide, t["bg"])
    add_rect(slide, 0, 0, SLIDE_W, Cm(0.55), t["accent"])

    company = meta.get("company", "")
    if company:
        add_text(slide, MARGIN, Cm(1.0), SLIDE_W - MARGIN * 2, Cm(0.8),
            company, t["font"], 11, t["secondary"], align=PP_ALIGN.RIGHT
        )

    headline = data.get("headline") or meta.get("title", "")
    add_text(slide, MARGIN, Cm(5.8), SLIDE_W - MARGIN * 2, Cm(5.5),
        headline, t["font"], 38, t["primary"], bold=True, wrap=True
    )

    sub = data.get("subheadline") or meta.get("subtitle", "")
    if sub:
        add_text(slide, MARGIN, Cm(12.0), SLIDE_W - MARGIN * 2, Cm(2.2),
            sub, t["font"], 17, t["secondary"], wrap=True
        )

    add_rect(slide, 0, SLIDE_H - Cm(1.9), SLIDE_W, Cm(1.9), t["surface"])

    date   = meta.get("date", "")
    author = meta.get("author", "")
    info   = f"{date}  ·  {author}" if date and author else date or author
    if info:
        add_text(slide, MARGIN, SLIDE_H - Cm(1.45), SLIDE_W - MARGIN * 2, Cm(1.0),
            info, t["font"], 11, t["secondary"]
        )


def build_agenda(slide, data, t, meta):
    set_bg(slide, t["bg"])
    title_bar(slide, data.get("headline", "목차"), t)

    items  = data.get("items", [])
    y0     = Cm(4.0)
    y_step = Cm(1.85)

    for i, item in enumerate(items):
        y = y0 + i * y_step
        add_text(slide, MARGIN, y, Cm(1.4), Cm(1.5),
            f"{i+1:02d}", t["font"], 18, t["accent"], bold=True
        )
        add_text(slide, MARGIN + Cm(1.7), y + Cm(0.1), Cm(26), Cm(1.4),
            item, t["font"], 17, t["primary"]
        )


def build_section_break(slide, data, t, meta):
    set_bg(slide, t["bg"])
    panel_w = Cm(11.0)
    add_rect(slide, 0, 0, panel_w, SLIDE_H, t["accent"])

    add_text(slide, Cm(0.5), Cm(5.8), panel_w - Cm(1.0), Cm(1.0),
        "SECTION", t["font"], 10, t["white"], align=PP_ALIGN.CENTER
    )

    num = data.get("section_number", "")
    if num:
        add_text(slide, Cm(0.5), Cm(6.6), panel_w - Cm(1.0), Cm(4.5),
            num, t["font"], 72, t["white"], bold=True, align=PP_ALIGN.CENTER
        )

    rx = panel_w + Cm(2.2)
    rw = SLIDE_W - panel_w - Cm(3.5)
    add_text(slide, rx, Cm(6.5), rw, Cm(3.5),
        data.get("headline", ""), t["font"], 30, t["primary"], bold=True, wrap=True
    )

    desc = data.get("description", "")
    if desc:
        add_text(slide, rx, Cm(10.8), rw, Cm(2.5),
            desc, t["font"], 14, t["secondary"], wrap=True
        )


def build_single(slide, data, t, meta):
    set_bg(slide, t["bg"])
    title_bar(slide, data.get("headline", ""), t)

    bullets = data.get("bullets", [])
    body    = data.get("body", "")
    y0      = Cm(3.9)
    avail_h = SLIDE_H - y0 - Cm(1.5)

    if bullets:
        add_bullets(slide, MARGIN, y0, SLIDE_W - MARGIN * 2, avail_h,
            bullets, t["font"], 16, t["primary"], t["secondary"]
        )
    elif body:
        add_text(slide, MARGIN, y0, SLIDE_W - MARGIN * 2, avail_h,
            body, t["font"], 16, t["primary"], wrap=True
        )

    source = data.get("source", "")
    if source:
        add_text(slide, MARGIN, SLIDE_H - Cm(1.4), SLIDE_W - MARGIN * 2, Cm(1.0),
            f"출처: {source}", t["font"], 9, t["secondary"]
        )


def build_two_column(slide, data, t, meta):
    set_bg(slide, t["bg"])
    title_bar(slide, data.get("headline", ""), t)

    col_gap = Cm(1.0)
    col_w   = (SLIDE_W - MARGIN * 2 - col_gap) / 2
    y0      = Cm(3.8)
    col_h   = SLIDE_H - y0 - Cm(1.5)

    add_rect(slide, MARGIN + col_w + Cm(0.42), y0, Cm(0.04), col_h, t["rule"])

    for side, sx in [("left", MARGIN), ("right", MARGIN + col_w + col_gap)]:
        col       = data.get(side, {})
        col_title = col.get("title", "")
        if col_title:
            add_text(slide, sx, y0, col_w, Cm(1.2),
                col_title, t["font"], 13, t["accent"], bold=True
            )
        dy = y0 + (Cm(1.4) if col_title else 0)
        dh = col_h - (Cm(1.4) if col_title else 0)

        if col.get("bullets"):
            add_bullets(slide, sx, dy, col_w, dh,
                col["bullets"], t["font"], 14, t["primary"], t["secondary"]
            )
        elif col.get("body"):
            add_text(slide, sx, dy, col_w, dh,
                col["body"], t["font"], 14, t["primary"], wrap=True
            )


def build_data_table(slide, data, t, meta):
    set_bg(slide, t["bg"])
    title_bar(slide, data.get("headline", ""), t)

    tdata   = data.get("table", {})
    headers = tdata.get("headers", [])
    rows    = tdata.get("rows", [])
    hi_col  = tdata.get("highlight_col", -1)

    if not headers and not rows:
        return

    n_cols = len(headers) or (len(rows[0]) if rows else 0)
    n_rows = len(rows) + 1
    row_h  = Cm(1.05)
    tbl    = slide.shapes.add_table(
        n_rows, n_cols, MARGIN, Cm(3.9),
        SLIDE_W - MARGIN * 2, int(row_h * n_rows)
    ).table

    col_w = int((SLIDE_W - MARGIN * 2) / n_cols)
    for i in range(n_cols):
        tbl.columns[i].width = col_w
    for i in range(n_rows):
        tbl.rows[i].height = int(row_h)

    font = t["font"]

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(t["accent"])
        tf = cell.text_frame
        tf.text = hdr
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = font
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = rgb(t["white"])

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if ci == hi_col:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(219, 234, 254)
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(t["surface"])
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(t["bg"])

            tf  = cell.text_frame
            tf.text = val
            p   = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name  = font
            run.font.size  = Pt(12)
            run.font.bold  = (ci == hi_col)
            run.font.color.rgb = rgb(t["primary"])

    source = data.get("source", "")
    if source:
        add_text(slide, MARGIN, SLIDE_H - Cm(1.4), SLIDE_W - MARGIN * 2, Cm(1.0),
            f"출처: {source}", font, 9, t["secondary"]
        )


def build_quote(slide, data, t, meta):
    set_bg(slide, t["bg"])
    add_rect(slide, 0, 0, SLIDE_W, Cm(0.35), t["accent"])
    add_rect(slide, MARGIN, Cm(5.5), Cm(0.35), Cm(7.0), t["accent"])

    add_text(slide, MARGIN + Cm(1.1), Cm(5.8),
        SLIDE_W - MARGIN * 2 - Cm(1.1), Cm(6.5),
        data.get("quote", ""), t["font"], 22, t["primary"], wrap=True
    )

    attribution = data.get("attribution", "")
    if attribution:
        add_text(slide, MARGIN, Cm(13.5), SLIDE_W - MARGIN * 2, Cm(1.2),
            "— " + attribution, t["font"], 12, t["secondary"],
            align=PP_ALIGN.RIGHT
        )


def build_closing(slide, data, t, meta):
    set_bg(slide, t["accent"])
    dark = tuple(max(v - 20, 0) for v in t["accent"])
    add_rect(slide, 0, SLIDE_H - Cm(2.0), SLIDE_W, Cm(2.0), dark)

    add_text(slide, MARGIN, Cm(6.0), SLIDE_W - MARGIN * 2, Cm(4.5),
        data.get("headline", "감사합니다"), t["font"], 46, t["white"],
        bold=True, align=PP_ALIGN.CENTER
    )

    contact = data.get("contact", "")
    if contact:
        add_text(slide, MARGIN, Cm(11.2), SLIDE_W - MARGIN * 2, Cm(1.5),
            contact, t["font"], 15, t["white"], align=PP_ALIGN.CENTER
        )

    note = data.get("note", "")
    if note:
        add_text(slide, MARGIN, Cm(13.2), SLIDE_W - MARGIN * 2, Cm(1.2),
            note, t["font"], 11, t["white"], align=PP_ALIGN.CENTER
        )

    company = meta.get("company", "")
    date    = meta.get("date", "")
    info    = f"{company}  ·  {date}" if company and date else company or date
    if info:
        add_text(slide, MARGIN, SLIDE_H - Cm(1.5), SLIDE_W - MARGIN * 2, Cm(1.0),
            info, t["font"], 10, t["white"], align=PP_ALIGN.CENTER
        )


# ─── Router ──────────────────────────────────────────────────────

BUILDERS = {
    "cover":         build_cover,
    "agenda":        build_agenda,
    "section_break": build_section_break,
    "single":        build_single,
    "two_column":    build_two_column,
    "data_table":    build_data_table,
    "quote":         build_quote,
    "closing":       build_closing,
}


def generate(spec_path: str, output_path: str, theme_override: str = None):
    with open(spec_path, encoding="utf-8") as f:
        spec = json.load(f)

    meta        = spec.get("meta", {})
    slides_data = spec.get("slides", [])
    theme_name  = theme_override or meta.get("theme", "consulting")

    if theme_name not in THEMES:
        print(f"Warning: unknown theme '{theme_name}', using 'consulting'", file=sys.stderr)
        theme_name = "consulting"

    t = dict(THEMES[theme_name])
    if meta.get("font"):
        t["font"] = meta["font"]

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for num, sdata in enumerate(slides_data, 1):
        layout  = sdata.get("layout", "single")
        builder = BUILDERS.get(layout)
        if not builder:
            print(f"Warning: unknown layout '{layout}', skipping slide {num}", file=sys.stderr)
            continue

        slide = prs.slides.add_slide(blank)
        builder(slide, sdata, t, meta)
        add_slide_number(slide, num, t)
        add_speaker_notes(slide, sdata.get("notes", ""))

    prs.save(output_path)
    print(f"  {output_path}  ({len(slides_data)} slides, theme: {theme_name})")


def main():
    parser = argparse.ArgumentParser(description="A-Team PPT Generator")
    parser.add_argument("spec",            help="JSON spec file path")
    parser.add_argument("--output", "-o", help="Output .pptx (default: <spec>.pptx)")
    parser.add_argument("--theme",
        choices=["consulting", "minimal", "executive"],
        help="Override theme"
    )
    args = parser.parse_args()

    spec_path = Path(args.spec)
    if not spec_path.exists():
        print(f"Error: {spec_path} not found", file=sys.stderr)
        sys.exit(1)

    out = args.output or str(spec_path.with_suffix(".pptx"))
    generate(str(spec_path), out, args.theme)


if __name__ == "__main__":
    main()
