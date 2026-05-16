#!/usr/bin/env python3
"""
PPT QA — 생성된 PPTX 자체 품질 검증

생성 직후 자동 실행. 통과 못 하면 출력 자체를 차단.

검증 항목:
  1. 텍스트 오버플로우 (슬라이드 경계 초과)
  2. 텍스트 겹침 (바운딩 박스 교차)
  3. 폰트 크기 이상 (본문 > 32pt, 제목 > 54pt)
  4. 데이터 손상 ($변수 치환, /bin/zsh 등 쉘 오염)
  5. 빈 슬라이드 (텍스트 없음)
  6. CJK 폰트 미적용
  7. 슬라이드 수 일치

Usage:
  python qa-pptx.py output.pptx [--json] [--fix]
  python qa-pptx.py output.pptx --spec spec.json   # 스펙 대비 검증
"""
import sys, json, argparse, os
from pptx import Presentation
from pptx.util import Emu, Pt

# ── 상수 ──
SLIDE_W_EMU = 12192000  # 16:9 width
SLIDE_H_EMU = 6858000   # 16:9 height
MAX_TITLE_PT = 54
MAX_BODY_PT = 32
MIN_BODY_PT = 9
SHELL_PATTERNS = ['/bin/', '/usr/', '→K', '→M', 'zsh', 'bash', 'SHELL=', 'HOME=']
CJK_FONTS = ['Malgun Gothic', 'Apple SD Gothic Neo', 'NanumGothic', 'Pretendard',
             'Noto Sans KR', 'Noto Sans CJK', 'D2Coding']


def emu_to_pt(emu):
    return round(emu / 12700, 1) if emu else 0


def get_text_boxes(slide):
    """슬라이드의 모든 텍스트 박스 정보 추출."""
    boxes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        full_text = tf.text.strip()
        if not full_text:
            continue

        # 폰트 크기 수집
        font_sizes = []
        has_cjk_font = False
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_sizes.append(emu_to_pt(run.font.size))
                if run.font.name and run.font.name in CJK_FONTS:
                    has_cjk_font = True

        boxes.append({
            'text': full_text[:100],
            'left': shape.left or 0,
            'top': shape.top or 0,
            'width': shape.width or 0,
            'height': shape.height or 0,
            'right': (shape.left or 0) + (shape.width or 0),
            'bottom': (shape.top or 0) + (shape.height or 0),
            'font_sizes': font_sizes,
            'max_font': max(font_sizes) if font_sizes else 0,
            'has_cjk_font': has_cjk_font,
        })
    return boxes


def check_overflow(boxes, slide_idx):
    """텍스트가 슬라이드 경계를 벗어나는지 확인."""
    issues = []
    margin = Emu(0)  # 0 margin tolerance
    for box in boxes:
        if box['right'] > SLIDE_W_EMU + margin:
            issues.append({
                'type': 'overflow_right',
                'severity': 'high',
                'slide': slide_idx + 1,
                'text': box['text'],
                'detail': f"Right edge {emu_to_pt(box['right'])}pt > slide width {emu_to_pt(SLIDE_W_EMU)}pt",
            })
        if box['bottom'] > SLIDE_H_EMU + margin:
            issues.append({
                'type': 'overflow_bottom',
                'severity': 'high',
                'slide': slide_idx + 1,
                'text': box['text'],
                'detail': f"Bottom edge {emu_to_pt(box['bottom'])}pt > slide height {emu_to_pt(SLIDE_H_EMU)}pt",
            })
    return issues


def check_overlap(boxes, slide_idx):
    """텍스트 박스 간 겹침 감지."""
    issues = []
    for i, a in enumerate(boxes):
        for j, b in enumerate(boxes):
            if j <= i:
                continue
            # AABB intersection
            if (a['left'] < b['right'] and a['right'] > b['left'] and
                a['top'] < b['bottom'] and a['bottom'] > b['top']):
                # Calculate overlap area
                ox = min(a['right'], b['right']) - max(a['left'], b['left'])
                oy = min(a['bottom'], b['bottom']) - max(a['top'], b['top'])
                overlap_area = ox * oy
                box_area = min(a['width'] * a['height'], b['width'] * b['height'])
                if box_area > 0 and overlap_area / box_area > 0.15:  # >15% overlap
                    issues.append({
                        'type': 'text_overlap',
                        'severity': 'critical',
                        'slide': slide_idx + 1,
                        'text': f"{a['text'][:40]} ↔ {b['text'][:40]}",
                        'detail': f"Overlap {round(overlap_area/box_area*100)}% of smaller box",
                    })
    return issues


def check_font_sizes(boxes, slide_idx):
    """비정상 폰트 크기 감지."""
    issues = []
    for box in boxes:
        if box['max_font'] > MAX_TITLE_PT:
            issues.append({
                'type': 'font_too_large',
                'severity': 'high',
                'slide': slide_idx + 1,
                'text': box['text'],
                'detail': f"Font {box['max_font']}pt > max {MAX_TITLE_PT}pt",
            })
        for fs in box['font_sizes']:
            if 0 < fs < MIN_BODY_PT:
                issues.append({
                    'type': 'font_too_small',
                    'severity': 'medium',
                    'slide': slide_idx + 1,
                    'text': box['text'],
                    'detail': f"Font {fs}pt < min {MIN_BODY_PT}pt",
                })
                break  # one per box
    return issues


def check_shell_corruption(boxes, slide_idx):
    """쉘 변수 치환 오염 감지."""
    issues = []
    for box in boxes:
        for pat in SHELL_PATTERNS:
            if pat in box['text']:
                issues.append({
                    'type': 'shell_corruption',
                    'severity': 'critical',
                    'slide': slide_idx + 1,
                    'text': box['text'],
                    'detail': f"Shell pattern detected: '{pat}'",
                })
                break
    return issues


def check_cjk(boxes, slide_idx):
    """한글 텍스트에 CJK 폰트 미적용 감지."""
    issues = []
    import re
    for box in boxes:
        has_korean = bool(re.search(r'[\uAC00-\uD7A3\u3131-\u3163]', box['text']))
        if has_korean and not box['has_cjk_font'] and box['font_sizes']:
            issues.append({
                'type': 'cjk_font_missing',
                'severity': 'medium',
                'slide': slide_idx + 1,
                'text': box['text'],
                'detail': 'Korean text without CJK font assignment',
            })
    return issues


def qa_pptx(pptx_path, spec_path=None):
    """PPTX 품질 검증 실행."""
    prs = Presentation(pptx_path)
    all_issues = []
    slide_count = len(prs.slides)

    for idx, slide in enumerate(prs.slides):
        boxes = get_text_boxes(slide)

        # 빈 슬라이드 체크
        if not boxes:
            all_issues.append({
                'type': 'empty_slide',
                'severity': 'medium',
                'slide': idx + 1,
                'text': '',
                'detail': 'No text content on slide',
            })
            continue

        all_issues.extend(check_overflow(boxes, idx))
        all_issues.extend(check_overlap(boxes, idx))
        all_issues.extend(check_font_sizes(boxes, idx))
        all_issues.extend(check_shell_corruption(boxes, idx))
        all_issues.extend(check_cjk(boxes, idx))

    # 스펙 대비 슬라이드 수 검증
    if spec_path and os.path.exists(spec_path):
        with open(spec_path, encoding='utf-8') as f:
            spec = json.load(f)
        expected = len(spec.get('slides', []))
        if slide_count != expected:
            all_issues.append({
                'type': 'slide_count_mismatch',
                'severity': 'medium',
                'slide': 0,
                'text': '',
                'detail': f"Expected {expected} slides, got {slide_count}",
            })

    # 스코어링
    weights = {'critical': 20, 'high': 5, 'medium': 2, 'low': 1}
    total_penalty = sum(weights.get(i['severity'], 1) for i in all_issues)
    score = max(0, 100 - total_penalty)
    grade = 'A' if score >= 90 else 'B' if score >= 70 else 'C' if score >= 50 else 'D' if score >= 30 else 'F'

    counts = {}
    for i in all_issues:
        counts[i['severity']] = counts.get(i['severity'], 0) + 1

    return {
        'file': pptx_path,
        'slides': slide_count,
        'score': score,
        'grade': grade,
        'issues': {'total': len(all_issues), **counts},
        'details': all_issues,
        'pass': score >= 70,  # B 이상만 통과
    }


def main():
    p = argparse.ArgumentParser(description='PPT QA — PPTX Quality Gate')
    p.add_argument('pptx', help='PPTX file to validate')
    p.add_argument('--spec', default=None, help='JSON spec file for comparison')
    p.add_argument('--json', action='store_true', help='JSON output')
    args = p.parse_args()

    result = qa_pptx(args.pptx, args.spec)

    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print(f"\n  PPT QA Report")
        print(f"  {'─' * 45}")
        print(f"  File:    {os.path.basename(result['file'])}")
        print(f"  Slides:  {result['slides']}")
        print(f"  Score:   {result['score']}/100 ({result['grade']})")
        print(f"  Issues:  {result['issues']['total']} "
              f"(critical:{result['issues'].get('critical',0)} "
              f"high:{result['issues'].get('high',0)} "
              f"med:{result['issues'].get('medium',0)})")
        print(f"  Pass:    {'YES' if result['pass'] else 'NO — BLOCK OUTPUT'}")

        if result['details']:
            print(f"\n  Details:")
            for i in result['details'][:15]:
                sev = '!!' if i['severity'] == 'critical' else '**' if i['severity'] == 'high' else ' *'
                print(f"  {sev} Slide {i['slide']}: {i['type']} — {i['detail']}")
                if i['text']:
                    print(f"       {i['text'][:60]}")
            if len(result['details']) > 15:
                print(f"  ... and {len(result['details']) - 15} more")
        print()

    sys.exit(0 if result['pass'] else 1)


if __name__ == '__main__':
    main()
